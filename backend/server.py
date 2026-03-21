from fastapi import FastAPI, APIRouter, HTTPException, Depends, Header, UploadFile, File
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict, EmailStr
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timezone, timedelta
import bcrypt
import jwt
from enum import Enum
import io
from twilio.rest import Client as TwilioClient
from twilio.base.exceptions import TwilioRestException

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# JWT Config
JWT_SECRET = os.environ.get('JWT_SECRET', 'hubspot-crm-secret-key-2024-very-secure')
JWT_ALGORITHM = 'HS256'
JWT_EXPIRATION_HOURS = 24

app = FastAPI(title="SalesHub Platform", version="2.0.0")
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ============= ENUMS =============

class RoleType(str, Enum):
    SUPER_ADMIN = "super_admin"
    ORG_ADMIN = "org_admin"
    MANAGER = "manager"
    SALES_REP = "sales_rep"
    VIEWER = "viewer"

class AssignmentMode(str, Enum):
    ROUND_ROBIN = "round_robin"
    TERRITORY = "territory"
    MANUAL = "manual"
    DEFAULT_AGENT = "default_agent"

class PipelineStage(str, Enum):
    NEW = "new"
    CONTACTED = "contacted"
    NO_ANSWER = "no_answer"
    INTERESTED = "interested"
    FOLLOW_UP = "follow_up"
    BOOKED = "booked"
    WON = "won"
    LOST = "lost"

class CustomerLifecycle(str, Enum):
    LEAD = "lead"
    AI_CONTACTED = "ai_contacted"
    INTERESTED = "interested"
    OPPORTUNITY = "opportunity"
    CUSTOMER = "customer"
    REPEAT_CUSTOMER = "repeat_customer"

class ActivityType(str, Enum):
    AI_CALL = "ai_call"
    EMAIL = "email"
    NOTE = "note"
    MEETING = "meeting"
    PIPELINE_CHANGE = "pipeline_change"
    STATUS_CHANGE = "status_change"
    ASSIGNMENT = "assignment"
    SYSTEM = "system"
    CONVERSION = "conversion"

class CallStatus(str, Enum):
    PENDING = "pending"
    IN_PROGRESS = "in_progress"
    NO_ANSWER = "no_answer"
    BUSY = "busy"
    COMPLETED = "completed"
    INTERESTED = "interested"
    FOLLOW_UP = "follow_up"
    FAILED = "failed"

class NotificationType(str, Enum):
    NEW_LEAD = "new_lead"
    LEAD_ASSIGNED = "lead_assigned"
    AI_CALL_COMPLETED = "ai_call_completed"
    DEAL_STAGE_CHANGED = "deal_stage_changed"
    CUSTOMER_REPLY = "customer_reply"
    TASK_DUE = "task_due"
    LEAD_CONVERTED = "lead_converted"
    SYSTEM = "system"

class Permission(str, Enum):
    # Organization
    MANAGE_ORGANIZATION = "manage_organization"
    VIEW_ORGANIZATION = "view_organization"
    # Users
    MANAGE_USERS = "manage_users"
    VIEW_USERS = "view_users"
    INVITE_USERS = "invite_users"
    # Leads
    MANAGE_ALL_LEADS = "manage_all_leads"
    MANAGE_OWN_LEADS = "manage_own_leads"
    VIEW_ALL_LEADS = "view_all_leads"
    VIEW_OWN_LEADS = "view_own_leads"
    # Deals
    MANAGE_ALL_DEALS = "manage_all_deals"
    MANAGE_OWN_DEALS = "manage_own_deals"
    VIEW_ALL_DEALS = "view_all_deals"
    VIEW_OWN_DEALS = "view_own_deals"
    # Contacts
    MANAGE_CONTACTS = "manage_contacts"
    VIEW_CONTACTS = "view_contacts"
    # Analytics
    VIEW_ANALYTICS = "view_analytics"
    VIEW_TEAM_ANALYTICS = "view_team_analytics"
    # Settings
    MANAGE_SETTINGS = "manage_settings"

# Role permissions mapping
ROLE_PERMISSIONS = {
    RoleType.SUPER_ADMIN: list(Permission),
    RoleType.ORG_ADMIN: [
        Permission.VIEW_ORGANIZATION, Permission.MANAGE_ORGANIZATION, Permission.MANAGE_USERS, Permission.VIEW_USERS,
        Permission.INVITE_USERS, Permission.MANAGE_ALL_LEADS, Permission.VIEW_ALL_LEADS,
        Permission.MANAGE_ALL_DEALS, Permission.VIEW_ALL_DEALS, Permission.MANAGE_CONTACTS,
        Permission.VIEW_CONTACTS, Permission.VIEW_ANALYTICS, Permission.VIEW_TEAM_ANALYTICS,
        Permission.MANAGE_SETTINGS
    ],
    RoleType.MANAGER: [
        Permission.VIEW_ORGANIZATION, Permission.VIEW_USERS, Permission.INVITE_USERS,
        Permission.MANAGE_ALL_LEADS, Permission.VIEW_ALL_LEADS, Permission.MANAGE_ALL_DEALS,
        Permission.VIEW_ALL_DEALS, Permission.MANAGE_CONTACTS, Permission.VIEW_CONTACTS,
        Permission.VIEW_ANALYTICS, Permission.VIEW_TEAM_ANALYTICS
    ],
    RoleType.SALES_REP: [
        Permission.VIEW_ORGANIZATION, Permission.MANAGE_OWN_LEADS, Permission.VIEW_OWN_LEADS,
        Permission.MANAGE_OWN_DEALS, Permission.VIEW_OWN_DEALS, Permission.MANAGE_CONTACTS,
        Permission.VIEW_CONTACTS, Permission.VIEW_ANALYTICS
    ],
    RoleType.VIEWER: [
        Permission.VIEW_ORGANIZATION, Permission.VIEW_OWN_LEADS, Permission.VIEW_OWN_DEALS,
        Permission.VIEW_CONTACTS, Permission.VIEW_ANALYTICS
    ]
}

# ============= MODELS =============

# Organization Models
class OrganizationCreate(BaseModel):
    name: str
    domain: Optional[str] = None
    industry: Optional[str] = None
    size: Optional[str] = None

class OrganizationUpdate(BaseModel):
    name: Optional[str] = None
    domain: Optional[str] = None
    industry: Optional[str] = None
    size: Optional[str] = None
    settings: Optional[Dict[str, Any]] = None

class OrganizationResponse(BaseModel):
    id: str
    name: str
    domain: Optional[str] = None
    industry: Optional[str] = None
    size: Optional[str] = None
    settings: Dict[str, Any] = {}
    created_at: str
    owner_id: str
    member_count: int = 0

# User Models
class UserCreate(BaseModel):
    email: EmailStr
    password: str
    name: str
    role: RoleType = RoleType.SALES_REP

class UserInvite(BaseModel):
    email: EmailStr
    name: str
    role: RoleType = RoleType.SALES_REP

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class UserUpdate(BaseModel):
    name: Optional[str] = None
    role: Optional[RoleType] = None
    is_active: Optional[bool] = None

class UserResponse(BaseModel):
    id: str
    email: str
    name: str
    role: RoleType
    organization_id: Optional[str] = None
    organization_name: Optional[str] = None
    permissions: List[str] = []
    is_active: bool = True
    created_at: str
    last_login: Optional[str] = None

class AuthResponse(BaseModel):
    token: str
    user: UserResponse

# Contact Models (HubSpot-style)
class ContactCreate(BaseModel):
    first_name: str
    last_name: Optional[str] = None
    email: Optional[str] = None
    phone: Optional[str] = None
    company: Optional[str] = None
    job_title: Optional[str] = None
    linkedin: Optional[str] = None
    address: Optional[str] = None
    city: Optional[str] = None
    state: Optional[str] = None
    postcode: Optional[str] = None
    country: Optional[str] = None
    notes: Optional[str] = None
    tags: List[str] = []
    custom_fields: Dict[str, Any] = {}
    is_public: Optional[bool] = False

class ContactUpdate(BaseModel):
    first_name: Optional[str] = None
    last_name: Optional[str] = None
    email: Optional[str] = None
    phone: Optional[str] = None
    company: Optional[str] = None
    job_title: Optional[str] = None
    linkedin: Optional[str] = None
    address: Optional[str] = None
    city: Optional[str] = None
    country: Optional[str] = None
    notes: Optional[str] = None
    tags: Optional[List[str]] = None
    custom_fields: Optional[Dict[str, Any]] = None

class ContactResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    first_name: str
    last_name: Optional[str] = None
    email: Optional[str] = None
    phone: Optional[str] = None
    company: Optional[str] = None
    job_title: Optional[str] = None
    linkedin: Optional[str] = None
    address: Optional[str] = None
    city: Optional[str] = None
    state: Optional[str] = None
    postcode: Optional[str] = None
    country: Optional[str] = None
    notes: Optional[str] = None
    tags: List[str] = []
    custom_fields: Dict[str, Any] = {}
    is_public: Optional[bool] = False
    organization_id: Optional[str] = None
    owner_id: str
    owner_name: Optional[str] = None
    created_at: str
    updated_at: str

# Lead Models
class LeadCreate(BaseModel):
    name: str  # Clinic Name / Company Name
    email: Optional[str] = None
    phone: Optional[str] = None  # Mobile Number
    company: Optional[str] = None
    title: Optional[str] = None  # Job Title
    linkedin: Optional[str] = None
    company_size: Optional[str] = None
    industry: Optional[str] = None
    source: Optional[str] = None
    notes: Optional[str] = None
    contact_id: Optional[str] = None
    address: Optional[str] = None
    postcode: Optional[str] = None
    city: Optional[str] = None
    state: Optional[str] = None
    country: Optional[str] = None
    is_public: Optional[bool] = False
    lifecycle_stage: Optional[str] = "lead"
    # New fields from PDF spec
    website: Optional[str] = None
    pic_name: Optional[str] = None  # Person in Charge
    office_number: Optional[str] = None
    fax_number: Optional[str] = None
    pipeline_status: Optional[str] = "new"

class LeadUpdate(BaseModel):
    name: Optional[str] = None
    email: Optional[str] = None
    phone: Optional[str] = None
    company: Optional[str] = None
    title: Optional[str] = None
    linkedin: Optional[str] = None
    company_size: Optional[str] = None
    industry: Optional[str] = None
    source: Optional[str] = None
    notes: Optional[str] = None
    status: Optional[str] = None
    pipeline_stage: Optional[str] = None
    assigned_to: Optional[str] = None
    address: Optional[str] = None
    postcode: Optional[str] = None
    city: Optional[str] = None
    state: Optional[str] = None
    country: Optional[str] = None
    is_public: Optional[bool] = None
    lifecycle_stage: Optional[str] = None
    # New fields
    website: Optional[str] = None
    pic_name: Optional[str] = None
    office_number: Optional[str] = None
    fax_number: Optional[str] = None
    pipeline_status: Optional[str] = None

class LeadResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    name: str
    email: Optional[str] = None
    phone: Optional[str] = None
    company: Optional[str] = None
    title: Optional[str] = None
    linkedin: Optional[str] = None
    company_size: Optional[str] = None
    industry: Optional[str] = None
    source: Optional[str] = None
    notes: Optional[str] = None
    status: str
    pipeline_stage: Optional[str] = "new"
    lifecycle_stage: Optional[str] = "lead"
    ai_score: int
    ai_insights: Optional[str] = None
    contact_id: Optional[str] = None
    address: Optional[str] = None
    postcode: Optional[str] = None
    city: Optional[str] = None
    state: Optional[str] = None
    country: Optional[str] = None
    is_public: Optional[bool] = False
    converted_to_client: Optional[bool] = False
    client_id: Optional[str] = None
    organization_id: Optional[str] = None
    owner_id: str
    owner_name: Optional[str] = None
    assigned_to: Optional[str] = None
    assigned_to_name: Optional[str] = None
    created_at: str
    updated_at: str
    # New fields
    website: Optional[str] = None
    pic_name: Optional[str] = None
    office_number: Optional[str] = None
    fax_number: Optional[str] = None
    pipeline_status: Optional[str] = "new"

# Deal Models
class DealCreate(BaseModel):
    title: str
    value: float
    lead_id: Optional[str] = None
    contact_id: Optional[str] = None
    company: Optional[str] = None
    contact_name: Optional[str] = None
    stage: str = "lead"
    expected_close_date: Optional[str] = None
    notes: Optional[str] = None
    probability: Optional[int] = None
    knowledge_base_content: Optional[str] = None  # AI knowledge base for calls
    # New: linked companies (list of lead/customer IDs)
    linked_company_ids: List[str] = []

class DealUpdate(BaseModel):
    title: Optional[str] = None
    value: Optional[float] = None
    stage: Optional[str] = None
    expected_close_date: Optional[str] = None
    notes: Optional[str] = None
    probability: Optional[int] = None
    assigned_to: Optional[str] = None
    linked_company_ids: Optional[List[str]] = None
    knowledge_base_content: Optional[str] = None  # AI knowledge base for calls

class LinkedCompany(BaseModel):
    id: str
    name: str
    pic_name: Optional[str] = None
    location: Optional[str] = None
    mobile: Optional[str] = None
    entity_type: str = "lead"  # lead or customer

class DealResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    title: str
    value: float
    lead_id: Optional[str] = None
    contact_id: Optional[str] = None
    company: Optional[str] = None
    contact_name: Optional[str] = None
    stage: str
    expected_close_date: Optional[str] = None
    notes: Optional[str] = None
    probability: Optional[int] = None
    knowledge_base_content: Optional[str] = None  # AI knowledge base for calls
    ai_health_score: int
    organization_id: Optional[str] = None  # Optional for legacy deals
    owner_id: Optional[str] = None  # Made optional for deals created without owner
    owner_name: Optional[str] = None
    assigned_to: Optional[str] = None
    assigned_to_name: Optional[str] = None
    created_at: str
    updated_at: str
    # New: linked companies
    linked_company_ids: List[str] = []
    linked_companies: List[Dict[str, Any]] = []
    linked_companies_count: int = 0

# Activity Models
class ActivityCreate(BaseModel):
    type: str  # call, email, meeting, note, task
    subject: str
    description: Optional[str] = None
    contact_id: Optional[str] = None
    lead_id: Optional[str] = None
    deal_id: Optional[str] = None
    scheduled_at: Optional[str] = None
    completed: bool = False

class ActivityResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    type: str
    subject: str
    description: Optional[str] = None
    contact_id: Optional[str] = None
    lead_id: Optional[str] = None
    deal_id: Optional[str] = None
    scheduled_at: Optional[str] = None
    completed: bool
    organization_id: Optional[str] = None  # Optional for legacy activities
    owner_id: str
    owner_name: Optional[str] = None
    created_at: str

# Notification Models
class NotificationCreate(BaseModel):
    type: str
    title: str
    message: str
    link: Optional[str] = None
    entity_type: Optional[str] = None  # lead, contact, deal
    entity_id: Optional[str] = None

class NotificationResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    type: str
    title: str
    message: str
    link: Optional[str] = None
    entity_type: Optional[str] = None
    entity_id: Optional[str] = None
    read: bool = False
    organization_id: str
    user_id: str
    created_at: str

# AI Call Models (Placeholder for future integration)
class AICallCreate(BaseModel):
    contact_id: Optional[str] = None
    lead_id: Optional[str] = None
    phone_number: str
    script_type: Optional[str] = "default"
    notes: Optional[str] = None

# Lead-Deal Linkage Models - for tracking per-lead pipeline status
class LeadDealLinkageCreate(BaseModel):
    lead_id: str
    deal_id: str
    pipeline_status: str = "lead"  # lead, qualified, proposal, negotiation, sales_closed, lost
    notes: Optional[str] = None

class LeadDealLinkageUpdate(BaseModel):
    pipeline_status: Optional[str] = None
    notes: Optional[str] = None

class LeadDealLinkageResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    lead_id: str
    deal_id: str
    pipeline_status: str
    notes: Optional[str] = None
    organization_id: str
    created_at: str
    updated_at: str
    # Enriched fields
    lead_name: Optional[str] = None
    lead_company: Optional[str] = None
    lead_phone: Optional[str] = None
    deal_title: Optional[str] = None
    deal_value: Optional[float] = None

class AICallResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    contact_id: Optional[str] = None
    lead_id: Optional[str] = None
    phone_number: str
    status: str  # pending, in_progress, completed, failed
    call_result: Optional[str] = None  # no_answer, busy, completed, interested, follow_up
    duration_seconds: Optional[int] = None
    transcription: Optional[str] = None
    ai_summary: Optional[str] = None
    recording_url: Optional[str] = None
    organization_id: str
    initiated_by: str
    created_at: str
    completed_at: Optional[str] = None

# Worklist Models
class WorklistItem(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    record_number: int
    customer_name: str
    customer_email: Optional[str] = None
    customer_phone: Optional[str] = None
    assigned_agent: Optional[str] = None
    assigned_agent_name: Optional[str] = None
    status: str
    lifecycle_stage: str
    registration_time: str
    last_activity: Optional[str] = None
    entity_type: str  # lead or contact
    entity_id: str

# Analytics Models
class AnalyticsResponse(BaseModel):
    total_leads: int
    total_deals: int
    total_contacts: int
    total_pipeline_value: float
    won_deals_value: float
    conversion_rate: float
    leads_by_status: Dict[str, int]
    leads_by_source: Dict[str, int] = {}
    deals_by_stage: Dict[str, int]
    monthly_revenue: List[Dict[str, Any]]
    team_performance: List[Dict[str, Any]] = []
    recent_activities: List[Dict[str, Any]] = []
    organization_stats: Dict[str, Any] = {}

# Assignment Settings Models
class TerritoryMapping(BaseModel):
    state: str
    city: Optional[str] = None
    agent_id: str

class AssignmentSettingsUpdate(BaseModel):
    mode: str  # round_robin, territory, manual, default_agent
    default_agent_id: Optional[str] = None
    territories: Optional[List[TerritoryMapping]] = None

class AssignmentSettingsResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    organization_id: str
    mode: str
    default_agent_id: Optional[str] = None
    default_agent_name: Optional[str] = None
    territories: List[Dict[str, Any]] = []
    round_robin_index: int = 0

# Client/Service Models
class ServiceCreate(BaseModel):
    name: str
    description: Optional[str] = None
    amount: float
    purchase_date: Optional[str] = None
    renewal_date: Optional[str] = None
    status: str = "active"  # active, expired, cancelled

class ClientCreate(BaseModel):
    lead_id: str
    services: List[ServiceCreate] = []
    notes: Optional[str] = None

class ClientResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    lead_id: str
    customer_id: Optional[str] = None  # Link to customer record
    customer_name: str
    customer_email: Optional[str] = None
    customer_phone: Optional[str] = None
    company: Optional[str] = None
    services: List[Dict[str, Any]] = []
    total_value: float = 0
    notes: Optional[str] = None
    organization_id: str
    converted_by: str
    converted_by_name: Optional[str] = None
    created_at: str
    updated_at: str

# Task Models
class PaymentStatus(str, Enum):
    PAID = "paid"
    PARTIALLY_PAID = "partially_paid"
    UNPAID = "unpaid"

class TaskCreate(BaseModel):
    title: str
    description: Optional[str] = None
    lead_id: Optional[str] = None
    client_id: Optional[str] = None
    deal_id: Optional[str] = None  # NEW: Link to deal
    assigned_to: Optional[str] = None
    due_date: Optional[str] = None
    payment_status: str = "unpaid"
    payment_amount: Optional[float] = None
    paid_amount: Optional[float] = 0
    calendar_event_id: Optional[str] = None
    priority: str = "medium"  # low, medium, high
    status: str = "pending"  # NEW: pending, in_progress, completed

class TaskUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    lead_id: Optional[str] = None
    client_id: Optional[str] = None
    deal_id: Optional[str] = None
    assigned_to: Optional[str] = None
    due_date: Optional[str] = None
    payment_status: Optional[str] = None
    payment_amount: Optional[float] = None
    paid_amount: Optional[float] = None
    status: Optional[str] = None
    priority: Optional[str] = None

class TaskResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    title: str
    description: Optional[str] = None
    lead_id: Optional[str] = None
    lead_name: Optional[str] = None
    client_id: Optional[str] = None
    client_name: Optional[str] = None
    deal_id: Optional[str] = None
    deal_name: Optional[str] = None
    deal_value: Optional[float] = None
    company_name: Optional[str] = None  # From lead or client
    pic_name: Optional[str] = None  # Person in charge
    assigned_to: Optional[str] = None
    assigned_to_name: Optional[str] = None
    due_date: Optional[str] = None
    payment_status: str = "unpaid"
    payment_amount: Optional[float] = None
    paid_amount: Optional[float] = 0
    status: str = "pending"  # pending, in_progress, completed
    priority: str = "medium"
    calendar_event_id: Optional[str] = None
    organization_id: str
    created_by: str
    created_at: str
    updated_at: str
    reg_time: Optional[str] = None  # Same as created_at, for display

# Organization Settings Models
class OrganizationSettingsUpdate(BaseModel):
    currency: Optional[str] = None  # USD, MYR, EUR, GBP, etc.
    currency_symbol: Optional[str] = None  # $, RM, €, £
    google_calendar_client_id: Optional[str] = None
    google_calendar_client_secret: Optional[str] = None
    google_calendar_enabled: Optional[bool] = None
    # Twilio WhatsApp Integration
    twilio_account_sid: Optional[str] = None
    twilio_auth_token: Optional[str] = None
    twilio_whatsapp_number: Optional[str] = None  # e.g., whatsapp:+14155238886
    twilio_enabled: Optional[bool] = None

class OrganizationSettingsResponse(BaseModel):
    model_config = ConfigDict(extra="ignore")
    organization_id: str
    currency: str = "USD"
    currency_symbol: str = "$"
    google_calendar_enabled: bool = False
    google_calendar_connected: bool = False
    # Twilio WhatsApp Integration
    twilio_enabled: bool = False
    twilio_connected: bool = False
    twilio_whatsapp_number: Optional[str] = None

# ============= AUTH HELPERS =============

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str, email: str, org_id: str = None) -> str:
    payload = {
        'user_id': user_id,
        'email': email,
        'organization_id': org_id,
        'exp': datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

async def get_current_user(authorization: str = Header(None)):
    if not authorization or not authorization.startswith('Bearer '):
        raise HTTPException(status_code=401, detail="Missing authorization header")
    
    token = authorization.split(' ')[1]
    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user = await db.users.find_one({'id': payload['user_id']}, {'_id': 0})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        
        # Add permissions based on role
        user['permissions'] = [p.value for p in ROLE_PERMISSIONS.get(RoleType(user['role']), [])]
        
        # Get organization name
        if user.get('organization_id'):
            org = await db.organizations.find_one({'id': user['organization_id']}, {'_id': 0, 'name': 1})
            user['organization_name'] = org['name'] if org else None
        
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

def check_permission(user: dict, permission: Permission):
    if permission.value not in user.get('permissions', []):
        raise HTTPException(status_code=403, detail=f"Permission denied: {permission.value}")

def get_data_filter(user: dict, permission_all: Permission, permission_own: Permission):
    """Get MongoDB filter based on user permissions"""
    base_filter = {'organization_id': user.get('organization_id')}
    
    if permission_all.value in user.get('permissions', []):
        return base_filter
    elif permission_own.value in user.get('permissions', []):
        return {**base_filter, 'owner_id': user['id']}
    else:
        raise HTTPException(status_code=403, detail="Permission denied")

# ============= AI HELPERS =============

async def calculate_lead_score(lead: dict) -> dict:
    score = 50
    insights = []
    
    if lead.get('email'):
        score += 10
        insights.append("Email provided")
    if lead.get('phone'):
        score += 10
        insights.append("Phone available")
    if lead.get('company'):
        score += 10
        insights.append("Company identified")
    if lead.get('title'):
        score += 5
        if any(role in lead['title'].lower() for role in ['ceo', 'cto', 'vp', 'director', 'head', 'manager']):
            score += 10
            insights.append("Decision maker")
    if lead.get('company_size') in ['51-200', '201-500', '500+']:
        score += 10
        insights.append("Mid-large company")
    
    return {"score": min(score, 100), "insights": " | ".join(insights) or "No specific insights"}

async def calculate_deal_health(deal: dict) -> int:
    stage_scores = {
        'lead': 30, 'qualified': 50, 'demo': 60, 'proposal': 75,
        'negotiation': 85, 'closed_won': 100, 'closed_lost': 0
    }
    return stage_scores.get(deal.get('stage', 'lead'), 50)

# ============= ORGANIZATION ROUTES =============

@api_router.post("/organizations", response_model=OrganizationResponse)
async def create_organization(org_data: OrganizationCreate, user: dict = Depends(get_current_user)):
    """Create a new organization (Super Admin only or first org for user)"""
    
    # Allow if super admin or user has no org
    if user.get('organization_id') and user['role'] != RoleType.SUPER_ADMIN.value:
        raise HTTPException(status_code=403, detail="You already belong to an organization")
    
    org_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    org_doc = {
        'id': org_id,
        'name': org_data.name,
        'domain': org_data.domain,
        'industry': org_data.industry,
        'size': org_data.size,
        'settings': {},
        'created_at': now,
        'owner_id': user['id']
    }
    
    await db.organizations.insert_one(org_doc)
    
    # Update user to be org admin of this org
    await db.users.update_one(
        {'id': user['id']},
        {'$set': {'organization_id': org_id, 'role': RoleType.ORG_ADMIN.value}}
    )
    
    org_doc.pop('_id', None)
    org_doc['member_count'] = 1
    return OrganizationResponse(**org_doc)

@api_router.get("/organizations", response_model=List[OrganizationResponse])
async def get_organizations(user: dict = Depends(get_current_user)):
    """Get organizations (Super Admin sees all, others see their own)"""
    
    if user['role'] == RoleType.SUPER_ADMIN.value:
        orgs = await db.organizations.find({}, {'_id': 0}).to_list(1000)
    else:
        if not user.get('organization_id'):
            return []
        orgs = await db.organizations.find({'id': user['organization_id']}, {'_id': 0}).to_list(1)
    
    # Add member counts
    for org in orgs:
        count = await db.users.count_documents({'organization_id': org['id']})
        org['member_count'] = count
    
    return [OrganizationResponse(**org) for org in orgs]

@api_router.get("/organizations/{org_id}", response_model=OrganizationResponse)
async def get_organization(org_id: str, user: dict = Depends(get_current_user)):
    if user['role'] != RoleType.SUPER_ADMIN.value and user.get('organization_id') != org_id:
        raise HTTPException(status_code=403, detail="Access denied")
    
    org = await db.organizations.find_one({'id': org_id}, {'_id': 0})
    if not org:
        raise HTTPException(status_code=404, detail="Organization not found")
    
    org['member_count'] = await db.users.count_documents({'organization_id': org_id})
    return OrganizationResponse(**org)

@api_router.put("/organizations/{org_id}", response_model=OrganizationResponse)
async def update_organization(org_id: str, org_data: OrganizationUpdate, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.MANAGE_SETTINGS)
    
    if user['role'] != RoleType.SUPER_ADMIN.value and user.get('organization_id') != org_id:
        raise HTTPException(status_code=403, detail="Access denied")
    
    update_data = {k: v for k, v in org_data.model_dump().items() if v is not None}
    if update_data:
        await db.organizations.update_one({'id': org_id}, {'$set': update_data})
    
    org = await db.organizations.find_one({'id': org_id}, {'_id': 0})
    org['member_count'] = await db.users.count_documents({'organization_id': org_id})
    return OrganizationResponse(**org)

# ============= USER MANAGEMENT ROUTES =============

@api_router.post("/auth/register", response_model=AuthResponse)
async def register(user_data: UserCreate):
    """Register a new user (creates new org or joins existing)"""
    
    existing = await db.users.find_one({'email': user_data.email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    user_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    # Check if this is the first user (make them super admin)
    user_count = await db.users.count_documents({})
    role = RoleType.SUPER_ADMIN if user_count == 0 else user_data.role
    
    user_doc = {
        'id': user_id,
        'email': user_data.email,
        'password': hash_password(user_data.password),
        'name': user_data.name,
        'role': role.value,
        'organization_id': None,
        'is_active': True,
        'created_at': now,
        'last_login': now
    }
    
    await db.users.insert_one(user_doc)
    
    token = create_token(user_id, user_data.email, None)
    permissions = [p.value for p in ROLE_PERMISSIONS.get(role, [])]
    
    return AuthResponse(
        token=token,
        user=UserResponse(
            id=user_id,
            email=user_data.email,
            name=user_data.name,
            role=role,
            organization_id=None,
            organization_name=None,
            permissions=permissions,
            is_active=True,
            created_at=now,
            last_login=now
        )
    )

@api_router.post("/auth/login", response_model=AuthResponse)
async def login(credentials: UserLogin):
    user = await db.users.find_one({'email': credentials.email}, {'_id': 0})
    if not user or not verify_password(credentials.password, user['password']):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    
    if not user.get('is_active', True):
        raise HTTPException(status_code=403, detail="Account is deactivated")
    
    # Update last login
    await db.users.update_one(
        {'id': user['id']},
        {'$set': {'last_login': datetime.now(timezone.utc).isoformat()}}
    )
    
    token = create_token(user['id'], user['email'], user.get('organization_id'))
    permissions = [p.value for p in ROLE_PERMISSIONS.get(RoleType(user['role']), [])]
    
    org_name = None
    if user.get('organization_id'):
        org = await db.organizations.find_one({'id': user['organization_id']}, {'_id': 0, 'name': 1})
        org_name = org['name'] if org else None
    
    return AuthResponse(
        token=token,
        user=UserResponse(
            id=user['id'],
            email=user['email'],
            name=user['name'],
            role=RoleType(user['role']),
            organization_id=user.get('organization_id'),
            organization_name=org_name,
            permissions=permissions,
            is_active=user.get('is_active', True),
            created_at=user['created_at'],
            last_login=user.get('last_login')
        )
    )

@api_router.get("/auth/me", response_model=UserResponse)
async def get_me(user: dict = Depends(get_current_user)):
    return UserResponse(
        id=user['id'],
        email=user['email'],
        name=user['name'],
        role=RoleType(user['role']),
        organization_id=user.get('organization_id'),
        organization_name=user.get('organization_name'),
        permissions=user.get('permissions', []),
        is_active=user.get('is_active', True),
        created_at=user['created_at'],
        last_login=user.get('last_login')
    )

@api_router.post("/users/invite", response_model=UserResponse)
async def invite_user(invite_data: UserInvite, user: dict = Depends(get_current_user)):
    """Invite a user to the organization"""
    check_permission(user, Permission.INVITE_USERS)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    existing = await db.users.find_one({'email': invite_data.email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    user_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    temp_password = str(uuid.uuid4())[:8]  # Temporary password
    
    user_doc = {
        'id': user_id,
        'email': invite_data.email,
        'password': hash_password(temp_password),
        'name': invite_data.name,
        'role': invite_data.role.value,
        'organization_id': user['organization_id'],
        'is_active': True,
        'created_at': now,
        'invited_by': user['id']
    }
    
    await db.users.insert_one(user_doc)
    
    permissions = [p.value for p in ROLE_PERMISSIONS.get(invite_data.role, [])]
    
    return UserResponse(
        id=user_id,
        email=invite_data.email,
        name=invite_data.name,
        role=invite_data.role,
        organization_id=user['organization_id'],
        organization_name=user.get('organization_name'),
        permissions=permissions,
        is_active=True,
        created_at=now
    )

class PaginatedUsersResponse(BaseModel):
    items: List[UserResponse]
    total: int
    page: int
    limit: int
    total_pages: int

@api_router.get("/users", response_model=PaginatedUsersResponse)
async def get_users(
    page: int = 1,
    limit: int = 10,
    user: dict = Depends(get_current_user)
):
    """Get users in the organization"""
    check_permission(user, Permission.VIEW_USERS)
    
    # Validate pagination params
    page = max(1, page)
    limit = min(max(1, limit), 100)  # Cap at 100 items per page
    skip = (page - 1) * limit
    
    query = {}
    if user['role'] != RoleType.SUPER_ADMIN.value:
        if not user.get('organization_id'):
            return PaginatedUsersResponse(items=[], total=0, page=1, limit=limit, total_pages=0)
        query['organization_id'] = user['organization_id']
    
    # Get total count
    total = await db.users.count_documents(query)
    total_pages = (total + limit - 1) // limit  # Ceiling division
    
    # Get paginated results
    users = await db.users.find(query, {'_id': 0, 'password': 0}).skip(skip).limit(limit).to_list(limit)
    
    result = []
    for u in users:
        org_name = None
        if u.get('organization_id'):
            org = await db.organizations.find_one({'id': u['organization_id']}, {'_id': 0, 'name': 1})
            org_name = org['name'] if org else None
        
        permissions = [p.value for p in ROLE_PERMISSIONS.get(RoleType(u['role']), [])]
        result.append(UserResponse(
            id=u['id'],
            email=u['email'],
            name=u['name'],
            role=RoleType(u['role']),
            organization_id=u.get('organization_id'),
            organization_name=org_name,
            permissions=permissions,
            is_active=u.get('is_active', True),
            created_at=u['created_at'],
            last_login=u.get('last_login')
        ))
    
    return PaginatedUsersResponse(
        items=result,
        total=total,
        page=page,
        limit=limit,
        total_pages=total_pages
    )

@api_router.put("/users/{user_id}", response_model=UserResponse)
async def update_user(user_id: str, user_data: UserUpdate, user: dict = Depends(get_current_user)):
    """Update user (role, status)"""
    check_permission(user, Permission.MANAGE_USERS)
    
    target_user = await db.users.find_one({'id': user_id}, {'_id': 0})
    if not target_user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Can only manage users in same org (unless super admin)
    if user['role'] != RoleType.SUPER_ADMIN.value:
        if target_user.get('organization_id') != user.get('organization_id'):
            raise HTTPException(status_code=403, detail="Cannot manage users outside your organization")
    
    update_data = {k: v.value if isinstance(v, RoleType) else v for k, v in user_data.model_dump().items() if v is not None}
    if update_data:
        await db.users.update_one({'id': user_id}, {'$set': update_data})
    
    updated_user = await db.users.find_one({'id': user_id}, {'_id': 0, 'password': 0})
    permissions = [p.value for p in ROLE_PERMISSIONS.get(RoleType(updated_user['role']), [])]
    
    return UserResponse(
        id=updated_user['id'],
        email=updated_user['email'],
        name=updated_user['name'],
        role=RoleType(updated_user['role']),
        organization_id=updated_user.get('organization_id'),
        permissions=permissions,
        is_active=updated_user.get('is_active', True),
        created_at=updated_user['created_at'],
        last_login=updated_user.get('last_login')
    )

@api_router.delete("/users/{user_id}")
async def delete_user(user_id: str, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.MANAGE_USERS)
    
    if user_id == user['id']:
        raise HTTPException(status_code=400, detail="Cannot delete yourself")
    
    target_user = await db.users.find_one({'id': user_id}, {'_id': 0})
    if not target_user:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user['role'] != RoleType.SUPER_ADMIN.value:
        if target_user.get('organization_id') != user.get('organization_id'):
            raise HTTPException(status_code=403, detail="Cannot delete users outside your organization")
    
    await db.users.delete_one({'id': user_id})
    return {"message": "User deleted"}

# ============= CONTACT ROUTES =============

@api_router.post("/contacts", response_model=ContactResponse)
async def create_contact(contact_data: ContactCreate, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    contact_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    contact_doc = {
        'id': contact_id,
        **contact_data.model_dump(),
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.contacts.insert_one(contact_doc)
    contact_doc.pop('_id', None)
    contact_doc['owner_name'] = user['name']
    
    return ContactResponse(**contact_doc)

class PaginatedContactsResponse(BaseModel):
    items: List[ContactResponse]
    total: int
    page: int
    limit: int
    total_pages: int

@api_router.get("/contacts", response_model=PaginatedContactsResponse)
async def get_contacts(
    search: Optional[str] = None, 
    page: int = 1, 
    limit: int = 10,
    user: dict = Depends(get_current_user)
):
    check_permission(user, Permission.VIEW_CONTACTS)
    
    if not user.get('organization_id'):
        return PaginatedContactsResponse(items=[], total=0, page=1, limit=limit, total_pages=0)
    
    # Validate pagination params
    page = max(1, page)
    limit = min(max(1, limit), 100)  # Cap at 100 items per page
    skip = (page - 1) * limit
    
    query = {'organization_id': user['organization_id']}
    if search:
        query['$or'] = [
            {'first_name': {'$regex': search, '$options': 'i'}},
            {'last_name': {'$regex': search, '$options': 'i'}},
            {'email': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}}
        ]
    
    # Get total count
    total = await db.contacts.count_documents(query)
    total_pages = (total + limit - 1) // limit  # Ceiling division
    
    # Get paginated results
    contacts = await db.contacts.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    # Add owner names
    for contact in contacts:
        owner = await db.users.find_one({'id': contact['owner_id']}, {'_id': 0, 'name': 1})
        contact['owner_name'] = owner['name'] if owner else None
    
    return PaginatedContactsResponse(
        items=[ContactResponse(**c) for c in contacts],
        total=total,
        page=page,
        limit=limit,
        total_pages=total_pages
    )

@api_router.get("/contacts/{contact_id}", response_model=ContactResponse)
async def get_contact(contact_id: str, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.VIEW_CONTACTS)
    
    contact = await db.contacts.find_one(
        {'id': contact_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    if not contact:
        raise HTTPException(status_code=404, detail="Contact not found")
    
    owner = await db.users.find_one({'id': contact['owner_id']}, {'_id': 0, 'name': 1})
    contact['owner_name'] = owner['name'] if owner else None
    
    return ContactResponse(**contact)

@api_router.put("/contacts/{contact_id}", response_model=ContactResponse)
async def update_contact(contact_id: str, contact_data: ContactUpdate, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    contact = await db.contacts.find_one({'id': contact_id, 'organization_id': user.get('organization_id')})
    if not contact:
        raise HTTPException(status_code=404, detail="Contact not found")
    
    update_data = {k: v for k, v in contact_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    await db.contacts.update_one({'id': contact_id}, {'$set': update_data})
    
    updated = await db.contacts.find_one({'id': contact_id}, {'_id': 0})
    owner = await db.users.find_one({'id': updated['owner_id']}, {'_id': 0, 'name': 1})
    updated['owner_name'] = owner['name'] if owner else None
    
    return ContactResponse(**updated)

@api_router.delete("/contacts/{contact_id}")
async def delete_contact(contact_id: str, user: dict = Depends(get_current_user)):
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    result = await db.contacts.delete_one({'id': contact_id, 'organization_id': user.get('organization_id')})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Contact not found")
    
    return {"message": "Contact deleted"}

@api_router.post("/contacts/import")
async def import_contacts(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """Import contacts from Excel file"""
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    try:
        import pandas as pd
        
        contents = await file.read()
        
        # Determine file type and read accordingly
        if file.filename.endswith('.csv'):
            df = pd.read_csv(io.BytesIO(contents))
        else:
            # Try to auto-detect header row by looking for 'Clinic Name' or 'Name'
            df_raw = pd.read_excel(io.BytesIO(contents), header=None)
            header_row = 0
            for i in range(min(5, len(df_raw))):
                row_values = [str(v).lower() if pd.notna(v) else '' for v in df_raw.iloc[i]]
                if any('clinic name' in v or v == 'name' for v in row_values):
                    header_row = i
                    break
            df = pd.read_excel(io.BytesIO(contents), header=header_row)
        
        logging.info(f"Excel columns found: {list(df.columns)}")
        
        # Column mapping from Excel to our schema (case-insensitive matching)
        column_mapping = {
            # Name variations
            'Clinic Name': 'first_name',
            'clinic name': 'first_name',
            'Name': 'first_name',
            'name': 'first_name',
            'First Name': 'first_name',
            'first_name': 'first_name',
            'FirstName': 'first_name',
            'Customer Name': 'first_name',
            'customer name': 'first_name',
            'Last Name': 'last_name',
            'last_name': 'last_name',
            'LastName': 'last_name',
            # Address variations
            'Address': 'address',
            'address': 'address',
            'Full Address': 'address',
            'Street Address': 'address',
            # Postcode variations
            'Postcode': 'postcode',
            'postcode': 'postcode',
            'Post Code': 'postcode',
            'Postal Code': 'postcode',
            'postal code': 'postcode',
            'Zip': 'postcode',
            'Zip Code': 'postcode',
            # City variations
            'City': 'city',
            'city': 'city',
            # State variations
            'State': 'state',
            'state': 'state',
            'Region': 'state',
            'Province': 'state',
            # Phone variations
            'Contact Number': 'phone',
            'contact number': 'phone',
            'Customer Number': 'phone',
            'customer number': 'phone',
            'Phone': 'phone',
            'phone': 'phone',
            'Phone Number': 'phone',
            'phone number': 'phone',
            'Mobile': 'phone',
            'mobile': 'phone',
            'Mobile Number': 'phone',
            'Tel': 'phone',
            'Telephone': 'phone',
            # Email variations
            'Email': 'email',
            'email': 'email',
            'Email Address': 'email',
            'E-mail': 'email',
            # Public flag variations
            'Is public': 'is_public',
            'is public': 'is_public',
            'Is Public': 'is_public',
            'Public': 'is_public',
            # Company variations
            'Company': 'company',
            'company': 'company',
            'Company Name': 'company',
            'Organization': 'company',
            # Job title variations
            'Job Title': 'job_title',
            'job_title': 'job_title',
            'Title': 'job_title',
            'Role': 'job_title',
            'Position': 'job_title',
        }
        
        # Rename columns
        df = df.rename(columns={k: v for k, v in column_mapping.items() if k in df.columns})
        logging.info(f"Columns after mapping: {list(df.columns)}")
        
        now = datetime.now(timezone.utc).isoformat()
        imported_count = 0
        
        for _, row in df.iterrows():
            contact_data = {
                'id': str(uuid.uuid4()),
                'first_name': str(row.get('first_name', '')).strip() if pd.notna(row.get('first_name')) else '',
                'last_name': str(row.get('last_name', '')).strip() if pd.notna(row.get('last_name')) else '',
                'email': str(row.get('email', '')).strip() if pd.notna(row.get('email')) else '',
                'phone': str(row.get('phone', '')).strip() if pd.notna(row.get('phone')) else '',
                'company': str(row.get('company', '')).strip() if pd.notna(row.get('company')) else '',
                'job_title': str(row.get('job_title', '')).strip() if pd.notna(row.get('job_title')) else '',
                'address': str(row.get('address', '')).strip() if pd.notna(row.get('address')) else '',
                'postcode': str(row.get('postcode', '')).strip() if pd.notna(row.get('postcode')) else '',
                'city': str(row.get('city', '')).strip() if pd.notna(row.get('city')) else '',
                'state': str(row.get('state', '')).strip() if pd.notna(row.get('state')) else '',
                'country': 'Malaysia',
                'is_public': str(row.get('is_public', '')).lower() in ['yes', 'true', '1'] if pd.notna(row.get('is_public')) else False,
                'linkedin': '',
                'notes': f"Imported from {file.filename}",
                'tags': [],
                'custom_fields': {},
                'organization_id': user['organization_id'],
                'owner_id': user['id'],
                'owner_name': user['name'],
                'created_at': now,
                'updated_at': now
            }
            
            # Skip rows without a name
            if not contact_data['first_name']:
                continue
            
            await db.contacts.insert_one(contact_data)
            imported_count += 1
        
        return {"imported": imported_count, "message": f"Successfully imported {imported_count} contacts"}
    
    except Exception as e:
        logging.error(f"Import error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to import file: {str(e)}")

# ============= CUSTOMERS ROUTES (Alias for Contacts) =============
# These routes use the 'contacts' collection but expose as '/customers' for frontend consistency

@api_router.get("/customers", response_model=PaginatedContactsResponse)
async def get_customers(
    search: Optional[str] = None, 
    page: int = 1,
    limit: int = 10,
    user: dict = Depends(get_current_user)
):
    """Get customers (contacts) for the organization"""
    query = get_data_filter(user, Permission.VIEW_CONTACTS, Permission.VIEW_CONTACTS)
    
    page = max(1, page)
    limit = min(max(1, limit), 100)
    skip = (page - 1) * limit
    
    if search:
        query['$or'] = [
            {'first_name': {'$regex': search, '$options': 'i'}},
            {'last_name': {'$regex': search, '$options': 'i'}},
            {'email': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}}
        ]
    
    total = await db.contacts.count_documents(query)
    total_pages = (total + limit - 1) // limit
    
    customers = await db.contacts.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    return PaginatedContactsResponse(
        items=[ContactResponse(**c) for c in customers],
        total=total,
        page=page,
        limit=limit,
        total_pages=total_pages
    )

@api_router.post("/customers", response_model=ContactResponse)
async def create_customer(contact_data: ContactCreate, user: dict = Depends(get_current_user)):
    """Create a new customer (contact)"""
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    contact_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    contact_doc = {
        'id': contact_id,
        **contact_data.model_dump(),
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.contacts.insert_one(contact_doc)
    contact_doc.pop('_id', None)
    
    return ContactResponse(**contact_doc)

@api_router.get("/customers/{customer_id}", response_model=ContactResponse)
async def get_customer(customer_id: str, user: dict = Depends(get_current_user)):
    """Get a specific customer (contact)"""
    query = get_data_filter(user, Permission.VIEW_CONTACTS, Permission.VIEW_CONTACTS)
    query['id'] = customer_id
    
    customer = await db.contacts.find_one(query, {'_id': 0})
    if not customer:
        raise HTTPException(status_code=404, detail="Customer not found")
    
    return ContactResponse(**customer)

@api_router.put("/customers/{customer_id}", response_model=ContactResponse)
async def update_customer(customer_id: str, contact_data: ContactUpdate, user: dict = Depends(get_current_user)):
    """Update a customer (contact)"""
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    update_data = {k: v for k, v in contact_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    result = await db.contacts.find_one_and_update(
        {'id': customer_id, 'organization_id': user.get('organization_id')},
        {'$set': update_data},
        return_document=True,
        projection={'_id': 0}
    )
    
    if not result:
        raise HTTPException(status_code=404, detail="Customer not found")
    
    return ContactResponse(**result)

@api_router.delete("/customers/{customer_id}")
async def delete_customer(customer_id: str, user: dict = Depends(get_current_user)):
    """Delete a customer (contact)"""
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    result = await db.contacts.delete_one(
        {'id': customer_id, 'organization_id': user.get('organization_id')}
    )
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Customer not found")
    
    return {"message": "Customer deleted successfully"}

@api_router.post("/customers/import")
async def import_customers(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """Import customers from Excel file (uses contacts import logic)"""
    return await import_contacts(file, user)

# ============= LEAD ROUTES =============

@api_router.post("/leads", response_model=LeadResponse)
async def create_lead(lead_data: LeadCreate, user: dict = Depends(get_current_user)):
    if Permission.MANAGE_ALL_LEADS.value not in user['permissions'] and Permission.MANAGE_OWN_LEADS.value not in user['permissions']:
        raise HTTPException(status_code=403, detail="Permission denied")
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    lead_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    lead_doc = {
        'id': lead_id,
        **lead_data.model_dump(),
        'status': 'new',
        'ai_score': 50,
        'ai_insights': None,
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'assigned_to': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    # Calculate AI score
    ai_result = await calculate_lead_score(lead_doc)
    lead_doc['ai_score'] = ai_result['score']
    lead_doc['ai_insights'] = ai_result['insights']
    
    await db.leads.insert_one(lead_doc)
    lead_doc.pop('_id', None)
    lead_doc['owner_name'] = user['name']
    lead_doc['assigned_to_name'] = user['name']
    
    return LeadResponse(**lead_doc)

class PaginatedLeadsResponse(BaseModel):
    items: List[LeadResponse]
    total: int
    page: int
    limit: int
    total_pages: int

@api_router.get("/leads", response_model=PaginatedLeadsResponse)
async def get_leads(
    status: Optional[str] = None, 
    pipeline_status: Optional[str] = None,
    search: Optional[str] = None,
    state: Optional[str] = None,
    country: Optional[str] = None,
    assigned_to: Optional[str] = None,
    page: int = 1,
    limit: int = 10,
    user: dict = Depends(get_current_user)
):
    query = get_data_filter(user, Permission.VIEW_ALL_LEADS, Permission.VIEW_OWN_LEADS)
    
    # Validate pagination params
    page = max(1, page)
    limit = min(max(1, limit), 100)  # Cap at 100 items per page
    skip = (page - 1) * limit
    
    if status:
        query['status'] = status
    if pipeline_status:
        query['pipeline_status'] = pipeline_status
    if state:
        query['state'] = state
    if country:
        query['country'] = country
    if assigned_to:
        query['assigned_to'] = assigned_to
    if search:
        query['$or'] = [
            {'name': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}},
            {'email': {'$regex': search, '$options': 'i'}},
            {'pic_name': {'$regex': search, '$options': 'i'}}
        ]
    
    # Get total count
    total = await db.leads.count_documents(query)
    total_pages = (total + limit - 1) // limit  # Ceiling division
    
    # Get paginated results
    leads = await db.leads.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    # Add owner and assigned names
    for lead in leads:
        owner = await db.users.find_one({'id': lead['owner_id']}, {'_id': 0, 'name': 1})
        lead['owner_name'] = owner['name'] if owner else None
        if lead.get('assigned_to'):
            assigned = await db.users.find_one({'id': lead['assigned_to']}, {'_id': 0, 'name': 1})
            lead['assigned_to_name'] = assigned['name'] if assigned else None
    
    return PaginatedLeadsResponse(
        items=[LeadResponse(**lead) for lead in leads],
        total=total,
        page=page,
        limit=limit,
        total_pages=total_pages
    )

@api_router.get("/leads/{lead_id}", response_model=LeadResponse)
async def get_lead(lead_id: str, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.VIEW_ALL_LEADS, Permission.VIEW_OWN_LEADS)
    query['id'] = lead_id
    
    lead = await db.leads.find_one(query, {'_id': 0})
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    owner = await db.users.find_one({'id': lead['owner_id']}, {'_id': 0, 'name': 1})
    lead['owner_name'] = owner['name'] if owner else None
    
    return LeadResponse(**lead)

@api_router.put("/leads/{lead_id}", response_model=LeadResponse)
async def update_lead(lead_id: str, lead_data: LeadUpdate, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.MANAGE_ALL_LEADS, Permission.MANAGE_OWN_LEADS)
    query['id'] = lead_id
    
    lead = await db.leads.find_one(query)
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    update_data = {k: v for k, v in lead_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    # Calculate AI score based on pipeline status
    pipeline_status = update_data.get('pipeline_status') or lead.get('pipeline_status', 'lead')
    pipeline_scores = {
        'lead': 20,
        'new': 20,
        'qualified': 40,
        'proposal': 60,
        'negotiation': 75,
        'closed': 100,
        'sales_closed': 100,
        'lost': 0
    }
    base_score = pipeline_scores.get(pipeline_status, 20)
    
    # Add bonus points for engagement
    bonus = 0
    if lead.get('email'): bonus += 5
    if lead.get('phone'): bonus += 5
    if lead.get('company'): bonus += 5
    if lead.get('notes'): bonus += 5
    
    update_data['ai_score'] = min(100, base_score + bonus)
    
    # Also recalculate if company info changes
    if any(k in update_data for k in ['company', 'title', 'company_size', 'email', 'phone']):
        lead.update(update_data)
        ai_result = await calculate_lead_score(lead)
        if ai_result['score'] > update_data['ai_score']:
            update_data['ai_score'] = ai_result['score']
        update_data['ai_insights'] = ai_result['insights']
    
    await db.leads.update_one({'id': lead_id}, {'$set': update_data})
    
    updated_lead = await db.leads.find_one({'id': lead_id}, {'_id': 0})
    owner = await db.users.find_one({'id': updated_lead['owner_id']}, {'_id': 0, 'name': 1})
    updated_lead['owner_name'] = owner['name'] if owner else None
    
    return LeadResponse(**updated_lead)

@api_router.delete("/leads/{lead_id}")
async def delete_lead(lead_id: str, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.MANAGE_ALL_LEADS, Permission.MANAGE_OWN_LEADS)
    query['id'] = lead_id
    
    result = await db.leads.delete_one(query)
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    return {"message": "Lead deleted"}

@api_router.get("/leads/{lead_id}/activities")
async def get_lead_activities(lead_id: str, user: dict = Depends(get_current_user)):
    """Get activities for a specific lead"""
    # Verify lead exists and user has access
    query = get_data_filter(user, Permission.VIEW_ALL_LEADS, Permission.VIEW_OWN_LEADS)
    query['id'] = lead_id
    
    lead = await db.leads.find_one(query, {'_id': 0})
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    # Get activities for this lead
    activities = await db.activities.find(
        {'entity_id': lead_id, 'entity_type': 'lead'},
        {'_id': 0}
    ).sort('created_at', -1).to_list(100)
    
    # Convert ObjectId if present
    for activity in activities:
        if '_id' in activity:
            del activity['_id']
    
    return {"activities": activities}

@api_router.post("/leads/{lead_id}/activities")
async def create_lead_activity(lead_id: str, activity_data: dict, user: dict = Depends(get_current_user)):
    """Log an activity for a specific lead"""
    # Verify lead exists and user has access
    query = get_data_filter(user, Permission.VIEW_ALL_LEADS, Permission.VIEW_OWN_LEADS)
    query['id'] = lead_id
    
    lead = await db.leads.find_one(query, {'_id': 0})
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    now = datetime.now(timezone.utc).isoformat()
    activity_doc = {
        'id': str(uuid.uuid4()),
        'entity_type': 'lead',
        'entity_id': lead_id,
        'type': activity_data.get('type', 'note'),
        'description': activity_data.get('description', ''),
        'notes': activity_data.get('notes', ''),
        'user_id': user['id'],
        'user_name': user['name'],
        'organization_id': user.get('organization_id'),
        'created_at': now
    }
    
    await db.activities.insert_one(activity_doc)
    if '_id' in activity_doc:
        del activity_doc['_id']
    
    # Update lead's updated_at
    await db.leads.update_one({'id': lead_id}, {'$set': {'updated_at': now}})
    
    return activity_doc

@api_router.post("/leads/import")
async def import_leads(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """Import leads from Excel file with auto-assignment"""
    if Permission.MANAGE_ALL_LEADS.value not in user['permissions'] and Permission.MANAGE_OWN_LEADS.value not in user['permissions']:
        raise HTTPException(status_code=403, detail="Permission denied")
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    try:
        import pandas as pd
        
        contents = await file.read()
        
        # Determine file type and read accordingly
        if file.filename.endswith('.csv'):
            df = pd.read_csv(io.BytesIO(contents))
        else:
            # Try to auto-detect header row by looking for 'Clinic Name' or 'Name'
            df_raw = pd.read_excel(io.BytesIO(contents), header=None)
            header_row = 0
            for i in range(min(5, len(df_raw))):
                row_values = [str(v).lower() if pd.notna(v) else '' for v in df_raw.iloc[i]]
                if any('clinic name' in v or v == 'name' for v in row_values):
                    header_row = i
                    break
            df = pd.read_excel(io.BytesIO(contents), header=header_row)
        
        logging.info(f"Excel columns found for leads: {list(df.columns)}")
        
        # Column mapping from Excel to our schema
        column_mapping = {
            'Clinic Name': 'name',
            'clinic name': 'name',
            'Name': 'name',
            'name': 'name',
            'Address': 'address',
            'address': 'address',
            'Postcode': 'postcode',
            'postcode': 'postcode',
            'City': 'city',
            'city': 'city',
            'State': 'state',
            'state': 'state',
            'Contact Number': 'phone',
            'contact number': 'phone',
            'Phone': 'phone',
            'phone': 'phone',
            'Email': 'email',
            'email': 'email',
            'Is public': 'is_public',
            'is public': 'is_public',
            'Company': 'company',
            'company': 'company',
            'Industry': 'industry',
            'industry': 'industry',
        }
        
        # Rename columns
        df = df.rename(columns={k: v for k, v in column_mapping.items() if k in df.columns})
        logging.info(f"Columns after mapping: {list(df.columns)}")
        
        # Get assignment settings
        assignment_settings = await db.assignment_settings.find_one(
            {'organization_id': user['organization_id']},
            {'_id': 0}
        )
        
        # Get sales reps for round-robin
        sales_reps = await db.users.find(
            {'organization_id': user['organization_id'], 'role': {'$in': ['sales_rep', 'manager', 'org_admin']}},
            {'_id': 0, 'id': 1, 'name': 1}
        ).to_list(100)
        
        round_robin_index = assignment_settings.get('round_robin_index', 0) if assignment_settings else 0
        
        now = datetime.now(timezone.utc).isoformat()
        imported_count = 0
        
        for _, row in df.iterrows():
            # Determine assigned agent based on settings
            assigned_to = None
            assigned_to_name = None
            
            if assignment_settings:
                mode = assignment_settings.get('mode', 'manual')
                
                if mode == 'default_agent' and assignment_settings.get('default_agent_id'):
                    assigned_to = assignment_settings['default_agent_id']
                    agent = await db.users.find_one({'id': assigned_to}, {'_id': 0, 'name': 1})
                    assigned_to_name = agent['name'] if agent else None
                    
                elif mode == 'round_robin' and sales_reps:
                    agent = sales_reps[round_robin_index % len(sales_reps)]
                    assigned_to = agent['id']
                    assigned_to_name = agent['name']
                    round_robin_index += 1
                    
                elif mode == 'territory':
                    territories = assignment_settings.get('territories', [])
                    lead_state = str(row.get('state', '')).strip().lower() if pd.notna(row.get('state')) else ''
                    lead_city = str(row.get('city', '')).strip().lower() if pd.notna(row.get('city')) else ''
                    
                    for territory in territories:
                        t_state = territory.get('state', '').lower()
                        t_city = territory.get('city', '').lower() if territory.get('city') else None
                        
                        if t_state == lead_state:
                            if t_city is None or t_city == lead_city:
                                assigned_to = territory.get('agent_id')
                                agent = await db.users.find_one({'id': assigned_to}, {'_id': 0, 'name': 1})
                                assigned_to_name = agent['name'] if agent else None
                                break
            
            lead_data = {
                'id': str(uuid.uuid4()),
                'name': str(row.get('name', '')).strip() if pd.notna(row.get('name')) else '',
                'email': str(row.get('email', '')).strip() if pd.notna(row.get('email')) else '',
                'phone': str(row.get('phone', '')).strip() if pd.notna(row.get('phone')) else '',
                'company': str(row.get('company', '')).strip() if pd.notna(row.get('company')) else '',
                'address': str(row.get('address', '')).strip() if pd.notna(row.get('address')) else '',
                'postcode': str(row.get('postcode', '')).strip() if pd.notna(row.get('postcode')) else '',
                'city': str(row.get('city', '')).strip() if pd.notna(row.get('city')) else '',
                'state': str(row.get('state', '')).strip() if pd.notna(row.get('state')) else '',
                'industry': str(row.get('industry', '')).strip() if pd.notna(row.get('industry')) else 'Healthcare',
                'is_public': str(row.get('is_public', '')).lower() in ['yes', 'true', '1'] if pd.notna(row.get('is_public')) else False,
                'title': '',
                'company_size': '',
                'source': 'import',
                'status': 'new',
                'pipeline_stage': 'new',
                'lifecycle_stage': 'lead',
                'notes': f"Imported from {file.filename}",
                'ai_score': 50,  # Default score for imports
                'organization_id': user['organization_id'],
                'owner_id': user['id'],
                'owner_name': user['name'],
                'assigned_to': assigned_to,
                'assigned_to_name': assigned_to_name,
                'created_at': now,
                'updated_at': now
            }
            
            # Skip rows without a name
            if not lead_data['name']:
                continue
            
            await db.leads.insert_one(lead_data)
            imported_count += 1
            
            # Create notification for assigned agent
            if assigned_to:
                await create_notification(
                    org_id=user['organization_id'],
                    user_id=assigned_to,
                    notification_type='lead_assigned',
                    title='New Lead Assigned',
                    message=f'Lead "{lead_data["name"]}" has been assigned to you',
                    link=f'/profile/lead/{lead_data["id"]}',
                    entity_type='lead',
                    entity_id=lead_data['id']
                )
        
        # Update round-robin index
        if assignment_settings and assignment_settings.get('mode') == 'round_robin':
            await db.assignment_settings.update_one(
                {'organization_id': user['organization_id']},
                {'$set': {'round_robin_index': round_robin_index}}
            )
        
        return {"imported": imported_count, "message": f"Successfully imported {imported_count} leads"}
    
    except Exception as e:
        logging.error(f"Import error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to import file: {str(e)}")

# ============= DEAL ROUTES =============

@api_router.post("/deals", response_model=DealResponse)
async def create_deal(deal_data: DealCreate, user: dict = Depends(get_current_user)):
    if Permission.MANAGE_ALL_DEALS.value not in user['permissions'] and Permission.MANAGE_OWN_DEALS.value not in user['permissions']:
        raise HTTPException(status_code=403, detail="Permission denied")
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    deal_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    deal_doc = {
        'id': deal_id,
        **deal_data.model_dump(),
        'ai_health_score': 50,
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'assigned_to': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    deal_doc['ai_health_score'] = await calculate_deal_health(deal_doc)
    
    await db.deals.insert_one(deal_doc)
    
    # Auto-create task for each linked company when deal is created
    linked_company_ids = deal_data.linked_company_ids or []
    if deal_data.lead_id and deal_data.lead_id not in linked_company_ids:
        linked_company_ids.append(deal_data.lead_id)
    
    for lead_id in linked_company_ids:
        lead = await db.leads.find_one({'id': lead_id}, {'_id': 0, 'name': 1, 'company': 1, 'pic_name': 1})
        if lead:
            task_id = str(uuid.uuid4())
            task_doc = {
                'id': task_id,
                'title': lead.get('company') or lead.get('name') or deal_data.title,
                'description': f'Follow up task for deal: {deal_data.title}',
                'lead_id': lead_id,
                'deal_id': deal_id,
                'assigned_to': user['id'],
                'due_date': None,
                'payment_status': 'unpaid',
                'payment_amount': None,
                'paid_amount': 0,
                'status': 'pending',
                'priority': 'medium',
                'calendar_event_id': None,
                'organization_id': user['organization_id'],
                'created_by': user['id'],
                'created_at': now,
                'updated_at': now
            }
            await db.tasks.insert_one(task_doc)
    
    deal_doc.pop('_id', None)
    deal_doc['owner_name'] = user['name']
    deal_doc['assigned_to_name'] = user['name']
    deal_doc['linked_companies'] = []
    deal_doc['linked_companies_count'] = len(deal_data.linked_company_ids)
    
    return DealResponse(**deal_doc)

@api_router.get("/deals", response_model=List[DealResponse])
async def get_deals(stage: Optional[str] = None, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.VIEW_ALL_DEALS, Permission.VIEW_OWN_DEALS)
    
    if stage:
        query['stage'] = stage
    
    deals = await db.deals.find(query, {'_id': 0}).sort('created_at', -1).to_list(1000)
    
    for deal in deals:
        owner = await db.users.find_one({'id': deal['owner_id']}, {'_id': 0, 'name': 1})
        deal['owner_name'] = owner['name'] if owner else None
        if deal.get('assigned_to'):
            assigned = await db.users.find_one({'id': deal['assigned_to']}, {'_id': 0, 'name': 1})
            deal['assigned_to_name'] = assigned['name'] if assigned else None
        
        # Populate linked companies
        linked_ids = deal.get('linked_company_ids', [])
        deal['linked_companies_count'] = len(linked_ids)
        linked_companies = []
        for lid in linked_ids[:3]:  # Only fetch first 3 for list view
            lead = await db.leads.find_one({'id': lid}, {'_id': 0, 'id': 1, 'name': 1, 'company': 1, 'pic_name': 1, 'phone': 1, 'city': 1, 'state': 1})
            if lead:
                linked_companies.append({
                    'id': lead['id'],
                    'name': lead.get('company') or lead.get('name'),
                    'pic_name': lead.get('pic_name') or lead.get('name'),
                    'location': f"{lead.get('city', '')}, {lead.get('state', '')}".strip(', '),
                    'mobile': lead.get('phone'),
                    'entity_type': 'lead'
                })
            else:
                # Check customers
                customer = await db.customers.find_one({'id': lid}, {'_id': 0, 'id': 1, 'first_name': 1, 'last_name': 1, 'company': 1, 'phone': 1, 'city': 1, 'state': 1})
                if customer:
                    linked_companies.append({
                        'id': customer['id'],
                        'name': customer.get('company') or f"{customer.get('first_name', '')} {customer.get('last_name', '')}",
                        'pic_name': f"{customer.get('first_name', '')} {customer.get('last_name', '')}",
                        'location': f"{customer.get('city', '')}, {customer.get('state', '')}".strip(', '),
                        'mobile': customer.get('phone'),
                        'entity_type': 'customer'
                    })
        deal['linked_companies'] = linked_companies
    
    return [DealResponse(**deal) for deal in deals]

@api_router.get("/deals/{deal_id}", response_model=DealResponse)
async def get_deal(deal_id: str, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.VIEW_ALL_DEALS, Permission.VIEW_OWN_DEALS)
    query['id'] = deal_id
    
    deal = await db.deals.find_one(query, {'_id': 0})
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    # Populate all linked companies for detail view
    linked_ids = deal.get('linked_company_ids', [])
    linked_companies = []
    for lid in linked_ids:
        lead = await db.leads.find_one({'id': lid}, {'_id': 0, 'id': 1, 'name': 1, 'company': 1, 'pic_name': 1, 'phone': 1, 'city': 1, 'state': 1})
        if lead:
            linked_companies.append({
                'id': lead['id'],
                'name': lead.get('company') or lead.get('name'),
                'pic_name': lead.get('pic_name') or lead.get('name'),
                'location': f"{lead.get('city', '')}, {lead.get('state', '')}".strip(', '),
                'mobile': lead.get('phone'),
                'entity_type': 'lead'
            })
        else:
            customer = await db.customers.find_one({'id': lid}, {'_id': 0, 'id': 1, 'first_name': 1, 'last_name': 1, 'company': 1, 'phone': 1, 'city': 1, 'state': 1})
            if customer:
                linked_companies.append({
                    'id': customer['id'],
                    'name': customer.get('company') or f"{customer.get('first_name', '')} {customer.get('last_name', '')}",
                    'pic_name': f"{customer.get('first_name', '')} {customer.get('last_name', '')}",
                    'location': f"{customer.get('city', '')}, {customer.get('state', '')}".strip(', '),
                    'mobile': customer.get('phone'),
                    'entity_type': 'customer'
                })
    
    deal['linked_companies'] = linked_companies
    deal['linked_companies_count'] = len(linked_ids)
    
    return DealResponse(**deal)

@api_router.put("/deals/{deal_id}", response_model=DealResponse)
async def update_deal(deal_id: str, deal_data: DealUpdate, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.MANAGE_ALL_DEALS, Permission.MANAGE_OWN_DEALS)
    query['id'] = deal_id
    
    deal = await db.deals.find_one(query)
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    update_data = {k: v for k, v in deal_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    if 'stage' in update_data:
        deal.update(update_data)
        update_data['ai_health_score'] = await calculate_deal_health(deal)
    
    await db.deals.update_one({'id': deal_id}, {'$set': update_data})
    
    updated_deal = await db.deals.find_one({'id': deal_id}, {'_id': 0})
    return DealResponse(**updated_deal)

@api_router.delete("/deals/{deal_id}")
async def delete_deal(deal_id: str, user: dict = Depends(get_current_user)):
    query = get_data_filter(user, Permission.MANAGE_ALL_DEALS, Permission.MANAGE_OWN_DEALS)
    query['id'] = deal_id
    
    result = await db.deals.delete_one(query)
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    return {"message": "Deal deleted"}

class KnowledgeBaseUpdate(BaseModel):
    knowledge_base_content: str

@api_router.put("/deals/{deal_id}/knowledge-base")
async def update_deal_knowledge_base(deal_id: str, kb_data: KnowledgeBaseUpdate, user: dict = Depends(get_current_user)):
    """Update the knowledge base content for a deal - used by AI calling"""
    query = get_data_filter(user, Permission.MANAGE_ALL_DEALS, Permission.MANAGE_OWN_DEALS)
    query['id'] = deal_id
    
    deal = await db.deals.find_one(query)
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    update_data = {
        'knowledge_base_content': kb_data.knowledge_base_content,
        'updated_at': datetime.now(timezone.utc).isoformat()
    }
    
    await db.deals.update_one({'id': deal_id}, {'$set': update_data})
    
    return {"success": True, "message": "Knowledge base updated", "deal_id": deal_id}

@api_router.get("/deals/{deal_id}/knowledge-base")
async def get_deal_knowledge_base(deal_id: str, user: dict = Depends(get_current_user)):
    """Get the knowledge base content for a deal"""
    query = get_data_filter(user, Permission.VIEW_ALL_DEALS, Permission.VIEW_OWN_DEALS)
    query['id'] = deal_id
    
    deal = await db.deals.find_one(query, {'_id': 0, 'id': 1, 'title': 1, 'knowledge_base_content': 1})
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    return {
        "deal_id": deal['id'],
        "title": deal['title'],
        "knowledge_base_content": deal.get('knowledge_base_content', '')
    }

def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """Extract text content from various file formats"""
    import io
    filename_lower = filename.lower()
    
    try:
        # PDF files
        if filename_lower.endswith('.pdf'):
            from PyPDF2 import PdfReader
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return '\n\n'.join(text_parts)
        
        # Word documents (.docx)
        elif filename_lower.endswith('.docx'):
            from docx import Document
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            return '\n\n'.join(text_parts)
        
        # PowerPoint files (.pptx)
        elif filename_lower.endswith('.pptx'):
            from pptx import Presentation
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            return '\n\n'.join(text_parts)
        
        # Excel files (.xlsx)
        elif filename_lower.endswith('.xlsx'):
            from openpyxl import load_workbook
            xlsx_file = io.BytesIO(file_content)
            wb = load_workbook(xlsx_file, read_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"--- Sheet: {sheet_name} ---"]
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    if row_text:
                        sheet_text.append(row_text)
                if len(sheet_text) > 1:
                    text_parts.append('\n'.join(sheet_text))
            return '\n\n'.join(text_parts)
        
        # Plain text files (.txt, .md, .csv)
        elif filename_lower.endswith(('.txt', '.md', '.csv', '.json')):
            return file_content.decode('utf-8')
        
        else:
            # Try to decode as text
            try:
                return file_content.decode('utf-8')
            except:
                raise ValueError(f"Unsupported file format: {filename}")
    
    except Exception as e:
        raise ValueError(f"Error extracting text from {filename}: {str(e)}")

@api_router.post("/deals/{deal_id}/knowledge-base/upload")
async def upload_deal_knowledge_base(
    deal_id: str, 
    file: UploadFile = File(...),
    user: dict = Depends(get_current_user)
):
    """Upload a knowledge base file (PDF, Word, PowerPoint, Excel, Text) for a deal"""
    query = get_data_filter(user, Permission.MANAGE_ALL_DEALS, Permission.MANAGE_OWN_DEALS)
    query['id'] = deal_id
    
    deal = await db.deals.find_one(query)
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    # Supported file types
    supported_extensions = ['.pdf', '.docx', '.pptx', '.xlsx', '.txt', '.md', '.csv', '.json']
    filename = file.filename.lower()
    
    if not any(filename.endswith(ext) for ext in supported_extensions):
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type. Supported formats: {', '.join(supported_extensions)}"
        )
    
    # Read file content
    file_content = await file.read()
    
    # Check file size (max 10MB)
    if len(file_content) > 10 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File too large. Maximum size is 10MB.")
    
    try:
        # Extract text from file
        extracted_text = extract_text_from_file(file_content, file.filename)
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract any text from the file.")
        
        # Update deal with knowledge base content
        update_data = {
            'knowledge_base_content': extracted_text,
            'knowledge_base_filename': file.filename,
            'knowledge_base_uploaded_at': datetime.now(timezone.utc).isoformat(),
            'updated_at': datetime.now(timezone.utc).isoformat()
        }
        
        await db.deals.update_one({'id': deal_id}, {'$set': update_data})
        
        return {
            "success": True, 
            "message": f"Knowledge base uploaded from {file.filename}",
            "deal_id": deal_id,
            "filename": file.filename,
            "extracted_length": len(extracted_text),
            "preview": extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
        }
    
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

# ============= LEAD-DEAL LINKAGE ROUTES =============
# These track the pipeline status per lead-deal combination

@api_router.post("/lead-deal-linkages")
async def create_or_update_linkage(linkage_data: LeadDealLinkageCreate, user: dict = Depends(get_current_user)):
    """Create or update a lead-deal linkage with pipeline status"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    # Check if linkage already exists
    existing = await db.lead_deal_linkages.find_one({
        'lead_id': linkage_data.lead_id,
        'deal_id': linkage_data.deal_id,
        'organization_id': user['organization_id']
    })
    
    now = datetime.now(timezone.utc).isoformat()
    
    if existing:
        # Update existing linkage
        update_data = {
            'pipeline_status': linkage_data.pipeline_status,
            'notes': linkage_data.notes,
            'updated_at': now
        }
        await db.lead_deal_linkages.update_one(
            {'id': existing['id']},
            {'$set': update_data}
        )
        linkage_id = existing['id']
    else:
        # Create new linkage
        linkage_id = str(uuid.uuid4())
        linkage_doc = {
            'id': linkage_id,
            'lead_id': linkage_data.lead_id,
            'deal_id': linkage_data.deal_id,
            'pipeline_status': linkage_data.pipeline_status,
            'notes': linkage_data.notes,
            'organization_id': user['organization_id'],
            'created_at': now,
            'updated_at': now
        }
        await db.lead_deal_linkages.insert_one(linkage_doc)
    
    # Fetch the linkage with enriched data
    linkage = await db.lead_deal_linkages.find_one({'id': linkage_id}, {'_id': 0})
    
    # Enrich with lead and deal info
    lead = await db.leads.find_one({'id': linkage['lead_id']}, {'_id': 0, 'name': 1, 'company': 1, 'phone': 1})
    deal = await db.deals.find_one({'id': linkage['deal_id']}, {'_id': 0, 'title': 1, 'value': 1})
    
    if lead:
        linkage['lead_name'] = lead.get('name')
        linkage['lead_company'] = lead.get('company')
        linkage['lead_phone'] = lead.get('phone')
    if deal:
        linkage['deal_title'] = deal.get('title')
        linkage['deal_value'] = deal.get('value')
    
    return linkage

@api_router.get("/lead-deal-linkages")
async def get_linkages(
    lead_id: Optional[str] = None,
    deal_id: Optional[str] = None,
    pipeline_status: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get all lead-deal linkages for the organization, optionally filtered"""
    if not user.get('organization_id'):
        return []
    
    query = {'organization_id': user['organization_id']}
    if lead_id:
        query['lead_id'] = lead_id
    if deal_id:
        query['deal_id'] = deal_id
    if pipeline_status:
        query['pipeline_status'] = pipeline_status
    
    linkages = await db.lead_deal_linkages.find(query, {'_id': 0}).sort('updated_at', -1).to_list(1000)
    
    # Enrich with lead and deal info
    for linkage in linkages:
        lead = await db.leads.find_one({'id': linkage['lead_id']}, {'_id': 0, 'name': 1, 'company': 1, 'phone': 1, 'state': 1})
        deal = await db.deals.find_one({'id': linkage['deal_id']}, {'_id': 0, 'title': 1, 'value': 1, 'expected_close_date': 1})
        
        if lead:
            linkage['lead_name'] = lead.get('name')
            linkage['lead_company'] = lead.get('company')
            linkage['lead_phone'] = lead.get('phone')
            linkage['lead_state'] = lead.get('state')
        if deal:
            linkage['deal_title'] = deal.get('title')
            linkage['deal_value'] = deal.get('value')
            linkage['deal_expected_close_date'] = deal.get('expected_close_date')
    
    return linkages

@api_router.put("/lead-deal-linkages/{linkage_id}")
async def update_linkage(linkage_id: str, linkage_data: LeadDealLinkageUpdate, user: dict = Depends(get_current_user)):
    """Update a specific lead-deal linkage"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    linkage = await db.lead_deal_linkages.find_one({
        'id': linkage_id,
        'organization_id': user['organization_id']
    })
    
    if not linkage:
        raise HTTPException(status_code=404, detail="Linkage not found")
    
    update_data = {k: v for k, v in linkage_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    await db.lead_deal_linkages.update_one({'id': linkage_id}, {'$set': update_data})
    
    updated = await db.lead_deal_linkages.find_one({'id': linkage_id}, {'_id': 0})
    return updated

@api_router.delete("/lead-deal-linkages/{linkage_id}")
async def delete_linkage(linkage_id: str, user: dict = Depends(get_current_user)):
    """Delete a lead-deal linkage"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    result = await db.lead_deal_linkages.delete_one({
        'id': linkage_id,
        'organization_id': user['organization_id']
    })
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Linkage not found")
    
    return {"message": "Linkage deleted"}

# ============= ACTIVITY ROUTES =============

@api_router.post("/activities", response_model=ActivityResponse)
async def create_activity(activity_data: ActivityCreate, user: dict = Depends(get_current_user)):
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    activity_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    activity_doc = {
        'id': activity_id,
        **activity_data.model_dump(),
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'created_at': now
    }
    
    await db.activities.insert_one(activity_doc)
    activity_doc.pop('_id', None)
    activity_doc['owner_name'] = user['name']
    
    return ActivityResponse(**activity_doc)

@api_router.get("/activities", response_model=List[ActivityResponse])
async def get_activities(
    contact_id: Optional[str] = None,
    lead_id: Optional[str] = None,
    deal_id: Optional[str] = None,
    entity_id: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    if not user.get('organization_id'):
        return []
    
    query = {'organization_id': user['organization_id']}
    
    # Support entity_id as a general filter (checks both lead_id and entity_id fields)
    if entity_id:
        query['$or'] = [
            {'lead_id': entity_id},
            {'entity_id': entity_id},
            {'contact_id': entity_id}
        ]
    if contact_id:
        query['contact_id'] = contact_id
    if lead_id:
        query['lead_id'] = lead_id
    if deal_id:
        query['deal_id'] = deal_id
    
    activities = await db.activities.find(query, {'_id': 0}).sort('created_at', -1).to_list(100)
    
    for activity in activities:
        # Handle missing owner_id gracefully
        owner_id = activity.get('owner_id') or activity.get('user_id')
        if owner_id:
            owner = await db.users.find_one({'id': owner_id}, {'_id': 0, 'name': 1})
            activity['owner_name'] = owner['name'] if owner else None
        else:
            activity['owner_name'] = activity.get('user_name')
        
        # Ensure owner_id exists for response model
        if 'owner_id' not in activity:
            activity['owner_id'] = owner_id or 'system'
    
    return [ActivityResponse(**a) for a in activities]

# ============= ANALYTICS ROUTES =============

@api_router.get("/analytics", response_model=AnalyticsResponse)
async def get_analytics(user: dict = Depends(get_current_user)):
    check_permission(user, Permission.VIEW_ANALYTICS)
    
    if not user.get('organization_id'):
        return AnalyticsResponse(
            total_leads=0, total_deals=0, total_contacts=0,
            total_pipeline_value=0, won_deals_value=0, conversion_rate=0,
            leads_by_status={}, leads_by_source={}, deals_by_stage={}, 
            monthly_revenue=[], team_performance=[], recent_activities=[], organization_stats={}
        )
    
    org_filter = {'organization_id': user['organization_id']}
    
    # If user can only see own data
    if Permission.VIEW_ALL_LEADS.value not in user['permissions']:
        org_filter['owner_id'] = user['id']
    
    # Get leads stats
    leads = await db.leads.find(org_filter, {'_id': 0}).to_list(10000)
    total_leads = len(leads)
    leads_by_status = {}
    leads_by_source = {}
    for lead in leads:
        # Use pipeline_status instead of status for consistency with Pipeline page
        status = lead.get('pipeline_status', lead.get('status', 'new'))
        leads_by_status[status] = leads_by_status.get(status, 0) + 1
        source = lead.get('source', 'other') or 'other'
        leads_by_source[source] = leads_by_source.get(source, 0) + 1
    
    # Get deals stats
    deals = await db.deals.find(org_filter, {'_id': 0}).to_list(10000)
    total_deals = len(deals)
    
    # Get lead-deal linkages for pipeline values (individual lead statuses)
    linkages = await db.lead_deal_linkages.find(org_filter, {'_id': 0}).to_list(10000)
    
    # Create a map of deal values
    deal_values = {d['id']: d.get('value', 0) for d in deals}
    
    # Calculate pipeline value from linkages (each linkage has its own status)
    total_pipeline_value = sum(
        deal_values.get(l.get('deal_id'), 0) 
        for l in linkages 
        if l.get('pipeline_status') not in ['sales_closed', 'lost']
    )
    
    # Calculate won deals value from linkages with 'sales_closed' status
    won_deals_value = sum(
        deal_values.get(l.get('deal_id'), 0) 
        for l in linkages 
        if l.get('pipeline_status') == 'sales_closed'
    )
    
    # If no linkages exist, fallback to deals-based calculation
    if not linkages:
        total_pipeline_value = sum(d.get('value', 0) for d in deals if d.get('stage') not in ['sales_closed', 'lost'])
        won_deals_value = sum(d.get('value', 0) for d in deals if d.get('stage') == 'sales_closed')
    
    # Count deals by stage (from linkages)
    deals_by_stage = {}
    for linkage in linkages:
        stage = linkage.get('pipeline_status', 'lead')
        deals_by_stage[stage] = deals_by_stage.get(stage, 0) + 1
    
    # If no linkages, fallback to deals
    if not linkages:
        for deal in deals:
            stage = deal.get('stage', 'lead')
            deals_by_stage[stage] = deals_by_stage.get(stage, 0) + 1
    
    # Conversion rate
    won_count = deals_by_stage.get('sales_closed', 0)
    total_closed = won_count + deals_by_stage.get('lost', 0)
    conversion_rate = round((won_count / total_closed * 100), 1) if total_closed > 0 else 0
    
    # Get contacts count
    total_contacts = await db.contacts.count_documents(org_filter)
    
    # Monthly revenue
    monthly_revenue = [
        {"month": "Jan", "revenue": won_deals_value * 0.1},
        {"month": "Feb", "revenue": won_deals_value * 0.15},
        {"month": "Mar", "revenue": won_deals_value * 0.12},
        {"month": "Apr", "revenue": won_deals_value * 0.18},
        {"month": "May", "revenue": won_deals_value * 0.2},
        {"month": "Jun", "revenue": won_deals_value * 0.25}
    ]
    
    # Recent activities (simulated based on recent leads/deals)
    recent_activities = []
    recent_leads = sorted(leads, key=lambda x: x.get('created_at', ''), reverse=True)[:3]
    recent_deals = sorted(deals, key=lambda x: x.get('created_at', ''), reverse=True)[:3]
    
    for lead in recent_leads:
        recent_activities.append({
            'type': 'lead_created',
            'title': f"New lead: {lead.get('name', 'Unknown')}",
            'description': f"From {lead.get('source', 'unknown')} source",
            'time': lead.get('created_at', '')[:10] if lead.get('created_at') else 'Recently',
            'icon': 'user'
        })
    
    for deal in recent_deals:
        recent_activities.append({
            'type': 'deal_created',
            'title': f"Deal: {deal.get('title', 'Unknown')}",
            'description': f"${deal.get('value', 0):,.0f} - {deal.get('stage', 'lead').replace('_', ' ')}",
            'time': deal.get('created_at', '')[:10] if deal.get('created_at') else 'Recently',
            'icon': 'briefcase'
        })
    
    # Sort activities by time
    recent_activities = sorted(recent_activities, key=lambda x: x.get('time', ''), reverse=True)[:5]
    
    # Organization stats for multi-tenancy view
    org_stats = {}
    if user.get('organization_id'):
        org = await db.organizations.find_one({'id': user['organization_id']}, {'_id': 0})
        if org:
            member_count = await db.users.count_documents({'organization_id': user['organization_id']})
            org_stats = {
                'name': org.get('name', 'Unknown'),
                'member_count': member_count,
                'domain': org.get('domain', ''),
                'industry': org.get('industry', ''),
                'created_at': org.get('created_at', '')[:10] if org.get('created_at') else ''
            }
    
    # Team performance (only for managers/admins)
    team_performance = []
    if Permission.VIEW_TEAM_ANALYTICS.value in user['permissions']:
        org_users = await db.users.find({'organization_id': user['organization_id']}, {'_id': 0}).to_list(100)
        for u in org_users:
            user_leads = await db.leads.count_documents({'owner_id': u['id'], 'organization_id': user['organization_id']})
            user_deals = await db.deals.count_documents({'owner_id': u['id'], 'organization_id': user['organization_id']})
            user_won = await db.deals.count_documents({'owner_id': u['id'], 'organization_id': user['organization_id'], 'stage': 'closed_won'})
            team_performance.append({
                'user_id': u['id'],
                'name': u['name'],
                'role': u['role'],
                'leads': user_leads,
                'deals': user_deals,
                'won_deals': user_won
            })
    
    return AnalyticsResponse(
        total_leads=total_leads,
        total_deals=total_deals,
        total_contacts=total_contacts,
        total_pipeline_value=total_pipeline_value,
        won_deals_value=won_deals_value,
        conversion_rate=round(conversion_rate, 1),
        leads_by_status=leads_by_status,
        leads_by_source=leads_by_source,
        deals_by_stage=deals_by_stage,
        monthly_revenue=monthly_revenue,
        team_performance=team_performance,
        recent_activities=recent_activities,
        organization_stats=org_stats
    )

# ============= ROLES INFO ROUTE =============

@api_router.get("/roles")
async def get_roles():
    """Get available roles and their permissions"""
    return {
        "roles": [
            {
                "id": role.value,
                "name": role.value.replace('_', ' ').title(),
                "permissions": [p.value for p in ROLE_PERMISSIONS.get(role, [])]
            }
            for role in RoleType
        ]
    }

# ============= ACTIVITY TIMELINE ROUTES =============

@api_router.get("/timeline/{entity_type}/{entity_id}")
async def get_entity_timeline(entity_type: str, entity_id: str, user: dict = Depends(get_current_user)):
    """Get activity timeline for a lead or contact"""
    if entity_type not in ['lead', 'contact']:
        raise HTTPException(status_code=400, detail="Invalid entity type")
    
    if not user.get('organization_id'):
        return []
    
    # Get activities for this entity
    query = {'organization_id': user['organization_id']}
    if entity_type == 'lead':
        query['lead_id'] = entity_id
    else:
        query['contact_id'] = entity_id
    
    activities = await db.activities.find(query, {'_id': 0}).sort('created_at', -1).to_list(100)
    
    # Add owner names - handle missing owner_id gracefully
    for activity in activities:
        owner_id = activity.get('owner_id') or activity.get('user_id')
        if owner_id:
            owner = await db.users.find_one({'id': owner_id}, {'_id': 0, 'name': 1})
            activity['owner_name'] = owner['name'] if owner else None
        else:
            activity['owner_name'] = activity.get('user_name')
    
    return activities

async def log_activity(
    org_id: str,
    user_id: str,
    activity_type: str,
    subject: str,
    description: str = None,
    lead_id: str = None,
    contact_id: str = None,
    deal_id: str = None
):
    """Helper function to log activities"""
    activity_doc = {
        'id': str(uuid.uuid4()),
        'type': activity_type,
        'subject': subject,
        'description': description,
        'lead_id': lead_id,
        'contact_id': contact_id,
        'deal_id': deal_id,
        'organization_id': org_id,
        'owner_id': user_id,
        'completed': True,
        'created_at': datetime.now(timezone.utc).isoformat()
    }
    await db.activities.insert_one(activity_doc)
    return activity_doc

# ============= NOTIFICATION ROUTES =============

@api_router.get("/notifications")
async def get_notifications(unread_only: bool = False, limit: int = 20, user: dict = Depends(get_current_user)):
    """Get notifications for current user"""
    if not user.get('organization_id'):
        return {"items": [], "unread_count": 0}
    
    query = {
        'organization_id': user['organization_id'],
        'user_id': user['id']
    }
    if unread_only:
        query['read'] = False
    
    notifications = await db.notifications.find(query, {'_id': 0}).sort('created_at', -1).limit(limit).to_list(limit)
    unread_count = await db.notifications.count_documents({
        'organization_id': user['organization_id'],
        'user_id': user['id'],
        'read': False
    })
    
    return {"items": notifications, "unread_count": unread_count}

@api_router.put("/notifications/{notification_id}/read")
async def mark_notification_read(notification_id: str, user: dict = Depends(get_current_user)):
    """Mark a notification as read"""
    result = await db.notifications.update_one(
        {'id': notification_id, 'user_id': user['id']},
        {'$set': {'read': True}}
    )
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Notification not found")
    return {"message": "Notification marked as read"}

@api_router.put("/notifications/read-all")
async def mark_all_notifications_read(user: dict = Depends(get_current_user)):
    """Mark all notifications as read"""
    await db.notifications.update_many(
        {'user_id': user['id'], 'read': False},
        {'$set': {'read': True}}
    )
    return {"message": "All notifications marked as read"}

async def create_notification(
    org_id: str,
    user_id: str,
    notification_type: str,
    title: str,
    message: str,
    link: str = None,
    entity_type: str = None,
    entity_id: str = None
):
    """Helper function to create notifications"""
    notification_doc = {
        'id': str(uuid.uuid4()),
        'type': notification_type,
        'title': title,
        'message': message,
        'link': link,
        'entity_type': entity_type,
        'entity_id': entity_id,
        'read': False,
        'organization_id': org_id,
        'user_id': user_id,
        'created_at': datetime.now(timezone.utc).isoformat()
    }
    await db.notifications.insert_one(notification_doc)
    return notification_doc

# ============= AI CALL ROUTES (PLACEHOLDER) =============

@api_router.post("/ai-calls")
async def initiate_ai_call(call_data: AICallCreate, user: dict = Depends(get_current_user)):
    """Initiate an AI call - PLACEHOLDER for future integration"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    call_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    call_doc = {
        'id': call_id,
        'contact_id': call_data.contact_id,
        'lead_id': call_data.lead_id,
        'phone_number': call_data.phone_number,
        'status': 'pending',
        'call_result': None,
        'duration_seconds': None,
        'transcription': None,
        'ai_summary': None,
        'recording_url': None,
        'organization_id': user['organization_id'],
        'initiated_by': user['id'],
        'created_at': now,
        'completed_at': None
    }
    
    await db.ai_calls.insert_one(call_doc)
    
    # Log activity
    await log_activity(
        org_id=user['organization_id'],
        user_id=user['id'],
        activity_type='ai_call',
        subject=f"AI call initiated to {call_data.phone_number}",
        description="Call queued for processing",
        lead_id=call_data.lead_id,
        contact_id=call_data.contact_id
    )
    
    call_doc.pop('_id', None)
    return {
        **call_doc,
        "message": "AI Call feature is a placeholder. Integration with voice provider required."
    }

@api_router.get("/ai-calls")
async def get_ai_calls(
    lead_id: Optional[str] = None,
    contact_id: Optional[str] = None,
    status: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get AI calls for organization"""
    if not user.get('organization_id'):
        return []
    
    query = {'organization_id': user['organization_id']}
    if lead_id:
        query['lead_id'] = lead_id
    if contact_id:
        query['contact_id'] = contact_id
    if status:
        query['status'] = status
    
    calls = await db.ai_calls.find(query, {'_id': 0}).sort('created_at', -1).to_list(100)
    return calls

@api_router.get("/ai-calls/{call_id}")
async def get_ai_call(call_id: str, user: dict = Depends(get_current_user)):
    """Get specific AI call details"""
    call = await db.ai_calls.find_one(
        {'id': call_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    if not call:
        raise HTTPException(status_code=404, detail="Call not found")
    return call

@api_router.get("/ai-calls/lead/{lead_id}")
async def get_lead_ai_calls(lead_id: str, user: dict = Depends(get_current_user)):
    """Get all AI calls for a specific lead"""
    if not user.get('organization_id'):
        return {"calls": []}
    
    calls = await db.ai_calls.find(
        {'lead_id': lead_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(50)
    
    return {"calls": calls}

class AiCallInitiate(BaseModel):
    lead_id: str
    deal_id: str
    agent_name: str
    phone: Optional[str] = None
    call_purpose: Optional[str] = "follow_up"  # follow_up, appointment, qualification, custom

@api_router.post("/ai-calls/initiate")
async def initiate_ai_call_endpoint(call_data: AiCallInitiate, user: dict = Depends(get_current_user)):
    """Initiate an AI call to a lead using ElevenLabs Conversational AI"""
    from ai_services import initiate_ai_call as ai_initiate_call
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    # Get lead details
    lead = await db.leads.find_one({'id': call_data.lead_id}, {'_id': 0})
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    # Get deal details
    deal = await db.deals.find_one({'id': call_data.deal_id}, {'_id': 0})
    if not deal:
        raise HTTPException(status_code=404, detail="Deal not found")
    
    # Get AI agent configuration
    agent = await db.ai_agents.find_one(
        {'organization_id': user['organization_id'], 'name': call_data.agent_name},
        {'_id': 0}
    )
    
    # Get knowledge base from deal if available
    knowledge_base = deal.get('knowledge_base_content')
    
    # Phone number to call
    phone_number = call_data.phone or lead.get('phone')
    if not phone_number:
        raise HTTPException(status_code=400, detail="No phone number available for this lead")
    
    # Initiate AI call via ElevenLabs
    result = await ai_initiate_call(
        phone_number=phone_number,
        lead_info=lead,
        deal_info=deal,
        agent_config=agent,
        knowledge_base=knowledge_base,
        call_purpose=call_data.call_purpose or "follow_up"
    )
    
    call_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    # Create AI call record
    call_doc = {
        'id': call_id,
        'lead_id': call_data.lead_id,
        'deal_id': call_data.deal_id,
        'agent_name': call_data.agent_name,
        'phone': phone_number,
        'lead_name': lead.get('name') or lead.get('company'),
        'deal_title': deal.get('title'),
        'status': 'initiated' if result.get('success') else 'failed',
        'direction': 'outbound',
        'source': 'AI Call',
        'duration': None,
        'summary': f'AI call initiated to discuss {deal.get("title")}. Agent: {call_data.agent_name}',
        'recording_url': None,
        'transcript': None,
        'call_purpose': call_data.call_purpose or "follow_up",
        'ai_script': result.get('call_data', {}).get('script'),
        'organization_id': user['organization_id'],
        'initiated_by': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.ai_calls.insert_one(call_doc)
    
    # Log activity
    await db.activities.insert_one({
        'id': str(uuid.uuid4()),
        'organization_id': user['organization_id'],
        'entity_type': 'lead',
        'entity_id': call_data.lead_id,
        'action': 'ai_call',
        'type': 'ai_call',
        'title': f'AI Call - {call_data.agent_name}',
        'description': f'AI call initiated to {lead.get("name") or lead.get("company")} to discuss {deal.get("title")}',
        'user_id': user['id'],
        'user_name': user['name'],
        'call_id': call_id,
        'metadata': {
            'agent_name': call_data.agent_name,
            'deal_title': deal.get('title'),
            'phone': phone_number,
            'call_purpose': call_data.call_purpose
        },
        'created_at': now
    })
    
    call_doc.pop('_id', None)
    return {"success": result.get('success', False), "call": call_doc, "message": result.get('message', f"AI call initiated with {call_data.agent_name}"), "next_steps": result.get('next_steps', [])}

# ============= AI AGENTS MANAGEMENT =============

class AiAgentCreate(BaseModel):
    name: str
    agent_id: str  # ElevenLabs Agent ID
    description: Optional[str] = None

@api_router.get("/ai-agents")
async def get_ai_agents(user: dict = Depends(get_current_user)):
    """Get all AI agents configured for the organization"""
    if not user.get('organization_id'):
        return {"agents": []}
    
    agents = await db.ai_agents.find(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(50)
    
    return {"agents": agents}

@api_router.post("/ai-agents")
async def create_ai_agent(agent_data: AiAgentCreate, user: dict = Depends(get_current_user)):
    """Create a new AI agent configuration"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    agent_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    agent_doc = {
        'id': agent_id,
        'name': agent_data.name,
        'agent_id': agent_data.agent_id,  # ElevenLabs Agent ID
        'description': agent_data.description or 'AI Voice Agent',
        'organization_id': user['organization_id'],
        'created_by': user['id'],
        'created_at': now
    }
    
    await db.ai_agents.insert_one(agent_doc)
    agent_doc.pop('_id', None)
    
    return {"success": True, "agent": agent_doc}

@api_router.delete("/ai-agents/{agent_id}")
async def delete_ai_agent(agent_id: str, user: dict = Depends(get_current_user)):
    """Delete an AI agent configuration"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    result = await db.ai_agents.delete_one({
        'id': agent_id,
        'organization_id': user.get('organization_id')
    })
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Agent not found")
    
    return {"success": True}

# ============= WORKLIST DASHBOARD ROUTES =============

@api_router.get("/worklist")
async def get_worklist(
    status: Optional[str] = None,
    lifecycle_stage: Optional[str] = None,
    assigned_to: Optional[str] = None,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    page: int = 1,
    limit: int = 20,
    user: dict = Depends(get_current_user)
):
    """Get worklist dashboard data - active leads requiring attention"""
    if not user.get('organization_id'):
        return {"items": [], "total": 0, "page": page, "limit": limit, "total_pages": 0}
    
    # Validate pagination
    page = max(1, page)
    limit = min(max(1, limit), 100)
    skip = (page - 1) * limit
    
    # Build query for leads (worklist focuses on leads)
    query = {'organization_id': user['organization_id']}
    
    # Apply filters
    if status:
        query['status'] = status
    if lifecycle_stage:
        query['lifecycle_stage'] = lifecycle_stage
    if assigned_to:
        query['assigned_to'] = assigned_to
    
    # Date filters
    if date_from:
        query['created_at'] = {'$gte': date_from}
    if date_to:
        if 'created_at' in query:
            query['created_at']['$lte'] = date_to
        else:
            query['created_at'] = {'$lte': date_to}
    
    # Get total count
    total = await db.leads.count_documents(query)
    total_pages = (total + limit - 1) // limit
    
    # Get paginated leads
    leads = await db.leads.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    # Transform to worklist items
    worklist_items = []
    record_number = skip + 1
    
    for lead in leads:
        assigned_name = None
        if lead.get('assigned_to'):
            assigned_user = await db.users.find_one({'id': lead['assigned_to']}, {'_id': 0, 'name': 1})
            assigned_name = assigned_user['name'] if assigned_user else None
        
        # Get last activity
        last_activity = await db.activities.find_one(
            {'lead_id': lead['id']},
            {'_id': 0, 'created_at': 1, 'subject': 1}
        )
        
        worklist_items.append({
            'id': f"WL-{record_number:04d}",
            'record_number': record_number,
            'customer_name': lead.get('name', ''),
            'customer_email': lead.get('email'),
            'customer_phone': lead.get('phone'),
            'assigned_agent': lead.get('assigned_to'),
            'assigned_agent_name': assigned_name,
            'status': lead.get('status', 'new'),
            'lifecycle_stage': lead.get('lifecycle_stage', 'lead'),
            'registration_time': lead.get('created_at', ''),
            'last_activity': last_activity.get('created_at') if last_activity else None,
            'entity_type': 'lead',
            'entity_id': lead['id']
        })
        record_number += 1
    
    return {
        "items": worklist_items,
        "total": total,
        "page": page,
        "limit": limit,
        "total_pages": total_pages
    }

# ============= CUSTOMER PROFILE ROUTES =============

@api_router.get("/profile/lead/{lead_id}")
async def get_lead_profile(lead_id: str, user: dict = Depends(get_current_user)):
    """Get comprehensive lead profile with all related data"""
    query = get_data_filter(user, Permission.VIEW_ALL_LEADS, Permission.VIEW_OWN_LEADS)
    query['id'] = lead_id
    
    lead = await db.leads.find_one(query, {'_id': 0})
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    # Get owner info
    owner = await db.users.find_one({'id': lead['owner_id']}, {'_id': 0, 'name': 1, 'email': 1})
    lead['owner_name'] = owner['name'] if owner else None
    lead['owner_email'] = owner['email'] if owner else None
    
    # Get assigned user info
    if lead.get('assigned_to'):
        assigned = await db.users.find_one({'id': lead['assigned_to']}, {'_id': 0, 'name': 1, 'email': 1})
        lead['assigned_to_name'] = assigned['name'] if assigned else None
    
    # Get activities/timeline
    activities = await db.activities.find(
        {'lead_id': lead_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(50)
    
    for activity in activities:
        act_owner = await db.users.find_one({'id': activity['owner_id']}, {'_id': 0, 'name': 1})
        activity['owner_name'] = act_owner['name'] if act_owner else None
    
    # Get AI calls
    ai_calls = await db.ai_calls.find(
        {'lead_id': lead_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(20)
    
    # Get related deals
    deals = await db.deals.find(
        {'lead_id': lead_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(10)
    
    return {
        "profile": lead,
        "activities": activities,
        "ai_calls": ai_calls,
        "deals": deals,
        "stats": {
            "total_activities": len(activities),
            "total_calls": len(ai_calls),
            "total_deals": len(deals)
        }
    }

@api_router.get("/profile/contact/{contact_id}")
async def get_contact_profile(contact_id: str, user: dict = Depends(get_current_user)):
    """Get comprehensive contact profile with all related data"""
    check_permission(user, Permission.VIEW_CONTACTS)
    
    contact = await db.contacts.find_one(
        {'id': contact_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    if not contact:
        raise HTTPException(status_code=404, detail="Contact not found")
    
    # Get owner info
    owner = await db.users.find_one({'id': contact['owner_id']}, {'_id': 0, 'name': 1, 'email': 1})
    contact['owner_name'] = owner['name'] if owner else None
    contact['owner_email'] = owner['email'] if owner else None
    
    # Get activities/timeline
    activities = await db.activities.find(
        {'contact_id': contact_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(50)
    
    for activity in activities:
        act_owner = await db.users.find_one({'id': activity['owner_id']}, {'_id': 0, 'name': 1})
        activity['owner_name'] = act_owner['name'] if act_owner else None
    
    # Get AI calls
    ai_calls = await db.ai_calls.find(
        {'contact_id': contact_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(20)
    
    # Get related deals
    deals = await db.deals.find(
        {'contact_id': contact_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(10)
    
    # Get related leads
    leads = await db.leads.find(
        {'contact_id': contact_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(10)
    
    return {
        "profile": contact,
        "activities": activities,
        "ai_calls": ai_calls,
        "deals": deals,
        "leads": leads,
        "stats": {
            "total_activities": len(activities),
            "total_calls": len(ai_calls),
            "total_deals": len(deals),
            "total_leads": len(leads)
        }
    }

@api_router.get("/profile/customer/{customer_id}")
async def get_customer_profile(customer_id: str, user: dict = Depends(get_current_user)):
    """Get comprehensive customer profile with all related data"""
    check_permission(user, Permission.VIEW_CONTACTS)
    
    customer = await db.customers.find_one(
        {'id': customer_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    if not customer:
        raise HTTPException(status_code=404, detail="Customer not found")
    
    # Add name field for consistency
    customer['name'] = f"{customer.get('first_name', '')} {customer.get('last_name', '')}".strip() or customer.get('company', 'Customer')
    
    # Get owner info
    if customer.get('owner_id'):
        owner = await db.users.find_one({'id': customer['owner_id']}, {'_id': 0, 'name': 1, 'email': 1})
        customer['owner_name'] = owner['name'] if owner else 'Unassigned'
        customer['owner_email'] = owner['email'] if owner else None
    else:
        customer['owner_name'] = 'Unassigned'
        customer['owner_email'] = None
    
    # Get activities/timeline
    activities = await db.activities.find(
        {'entity_id': customer_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(50)
    
    for activity in activities:
        if activity.get('user_id'):
            act_owner = await db.users.find_one({'id': activity['user_id']}, {'_id': 0, 'name': 1})
            activity['owner_name'] = act_owner['name'] if act_owner else activity.get('user_name', 'System')
        else:
            activity['owner_name'] = activity.get('user_name', 'System')
    
    # Get AI calls
    ai_calls = await db.ai_calls.find(
        {'customer_id': customer_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).sort('created_at', -1).to_list(20)
    
    # Get related deals
    deals = await db.deals.find(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(100)
    
    # Filter deals that have this customer's lead_id linked
    customer_deals = []
    if customer.get('lead_id'):
        customer_deals = [d for d in deals if customer['lead_id'] in (d.get('linked_company_ids') or [])]
    
    return {
        "profile": customer,
        "activities": activities,
        "ai_calls": ai_calls,
        "deals": customer_deals,
        "leads": [],
        "stats": {
            "total_activities": len(activities),
            "total_calls": len(ai_calls),
            "total_deals": len(customer_deals),
            "total_leads": 0
        }
    }

# ============= ADVANCED FILTER OPTIONS =============

@api_router.get("/filter-options")
async def get_filter_options(user: dict = Depends(get_current_user)):
    """Get available filter options for dropdowns"""
    if not user.get('organization_id'):
        return {}
    
    org_id = user['organization_id']
    
    # Get unique industries from leads
    industries = await db.leads.distinct('industry', {'organization_id': org_id})
    industries = [i for i in industries if i]
    
    # Get unique sources
    sources = await db.leads.distinct('source', {'organization_id': org_id})
    sources = [s for s in sources if s]
    
    # Get unique statuses
    statuses = ['new', 'contacted', 'qualified', 'lost']
    
    # Get lifecycle stages
    lifecycle_stages = [stage.value for stage in CustomerLifecycle]
    
    # Get users for assignment filter
    users = await db.users.find(
        {'organization_id': org_id},
        {'_id': 0, 'id': 1, 'name': 1}
    ).to_list(100)
    
    # Pipeline stages - call-based
    pipeline_stages = [stage.value for stage in PipelineStage]
    
    return {
        "industries": industries,
        "sources": sources,
        "statuses": statuses,
        "lifecycle_stages": lifecycle_stages,
        "users": users,
        "pipeline_stages": pipeline_stages,
        "date_presets": [
            {"label": "Today", "value": "today"},
            {"label": "Yesterday", "value": "yesterday"},
            {"label": "Last 7 Days", "value": "last_7_days"},
            {"label": "Last 30 Days", "value": "last_30_days"},
            {"label": "This Month", "value": "this_month"},
            {"label": "Last Month", "value": "last_month"},
            {"label": "Custom Range", "value": "custom"}
        ]
    }

# ============= ASSIGNMENT SETTINGS ROUTES =============

@api_router.get("/assignment-settings")
async def get_assignment_settings(user: dict = Depends(get_current_user)):
    """Get assignment settings for organization"""
    check_permission(user, Permission.VIEW_ORGANIZATION)
    
    if not user.get('organization_id'):
        return {"mode": "manual", "territories": [], "default_agent_id": None}
    
    settings = await db.assignment_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not settings:
        return {
            "organization_id": user['organization_id'],
            "mode": "manual",
            "default_agent_id": None,
            "default_agent_name": None,
            "territories": [],
            "round_robin_index": 0
        }
    
    # Get default agent name
    if settings.get('default_agent_id'):
        agent = await db.users.find_one({'id': settings['default_agent_id']}, {'_id': 0, 'name': 1})
        settings['default_agent_name'] = agent['name'] if agent else None
    
    return settings

@api_router.put("/assignment-settings")
async def update_assignment_settings(settings: AssignmentSettingsUpdate, user: dict = Depends(get_current_user)):
    """Update assignment settings for organization"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    now = datetime.now(timezone.utc).isoformat()
    
    settings_doc = {
        'organization_id': user['organization_id'],
        'mode': settings.mode,
        'default_agent_id': settings.default_agent_id,
        'territories': [t.model_dump() for t in (settings.territories or [])],
        'round_robin_index': 0,
        'updated_at': now
    }
    
    await db.assignment_settings.update_one(
        {'organization_id': user['organization_id']},
        {'$set': settings_doc},
        upsert=True
    )
    
    return {"message": "Assignment settings updated", "settings": settings_doc}

# ============= CLIENT ROUTES =============

@api_router.get("/clients")
async def get_clients(
    page: int = 1,
    limit: int = 10,
    search: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get all clients for organization"""
    check_permission(user, Permission.VIEW_CONTACTS)
    
    if not user.get('organization_id'):
        return {"items": [], "total": 0, "page": 1, "limit": limit, "total_pages": 0}
    
    page = max(1, page)
    limit = min(max(1, limit), 100)
    skip = (page - 1) * limit
    
    query = {'organization_id': user['organization_id']}
    
    if search:
        query['$or'] = [
            {'customer_name': {'$regex': search, '$options': 'i'}},
            {'customer_email': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}}
        ]
    
    total = await db.clients.count_documents(query)
    total_pages = (total + limit - 1) // limit
    
    clients = await db.clients.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    # Get converter names
    for client in clients:
        converter = await db.users.find_one({'id': client['converted_by']}, {'_id': 0, 'name': 1})
        client['converted_by_name'] = converter['name'] if converter else None
    
    return {
        "items": clients,
        "total": total,
        "page": page,
        "limit": limit,
        "total_pages": total_pages
    }

@api_router.get("/clients/{client_id}")
async def get_client(client_id: str, user: dict = Depends(get_current_user)):
    """Get client details"""
    check_permission(user, Permission.VIEW_CONTACTS)
    
    client = await db.clients.find_one(
        {'id': client_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    
    if not client:
        raise HTTPException(status_code=404, detail="Client not found")
    
    # Get converter name
    converter = await db.users.find_one({'id': client['converted_by']}, {'_id': 0, 'name': 1})
    client['converted_by_name'] = converter['name'] if converter else None
    
    # Get original lead
    lead = await db.leads.find_one({'id': client['lead_id']}, {'_id': 0})
    
    # Get activities
    activities = await db.activities.find(
        {'$or': [{'lead_id': client['lead_id']}, {'client_id': client_id}]},
        {'_id': 0}
    ).sort('created_at', -1).to_list(50)
    
    return {
        "client": client,
        "lead": lead,
        "activities": activities
    }

@api_router.post("/leads/{lead_id}/convert")
async def convert_lead_to_client(lead_id: str, client_data: ClientCreate, user: dict = Depends(get_current_user)):
    """Convert a lead to a client after successful sale"""
    if Permission.MANAGE_ALL_LEADS.value not in user['permissions'] and Permission.MANAGE_OWN_LEADS.value not in user['permissions']:
        raise HTTPException(status_code=403, detail="Permission denied")
    
    # Get the lead
    lead = await db.leads.find_one(
        {'id': lead_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    now = datetime.now(timezone.utc).isoformat()
    client_id = str(uuid.uuid4())
    
    # Calculate total value from services
    total_value = sum(s.amount for s in client_data.services)
    
    # Create client record
    client_doc = {
        'id': client_id,
        'lead_id': lead_id,
        'customer_id': None,  # Will be set if we create a customer record
        'customer_name': lead.get('name', ''),
        'customer_email': lead.get('email'),
        'customer_phone': lead.get('phone'),
        'company': lead.get('company'),
        'services': [s.model_dump() for s in client_data.services],
        'total_value': total_value,
        'notes': client_data.notes,
        'organization_id': user['organization_id'],
        'converted_by': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.clients.insert_one(client_doc)
    
    # Update lead status and lifecycle
    await db.leads.update_one(
        {'id': lead_id},
        {'$set': {
            'status': 'qualified',
            'pipeline_stage': 'won',
            'lifecycle_stage': 'customer',
            'converted_to_client': True,
            'client_id': client_id,
            'updated_at': now
        }}
    )
    
    # Create customer record in customers collection
    customer_id = str(uuid.uuid4())
    customer_doc = {
        'id': customer_id,
        'first_name': lead.get('name', '').split()[0] if lead.get('name') else '',
        'last_name': ' '.join(lead.get('name', '').split()[1:]) if lead.get('name') else '',
        'email': lead.get('email'),
        'phone': lead.get('phone'),
        'company': lead.get('company'),
        'job_title': lead.get('title'),
        'address': lead.get('address'),
        'city': lead.get('city'),
        'state': lead.get('state'),
        'postcode': lead.get('postcode'),
        'is_client': True,
        'client_id': client_id,
        'lead_id': lead_id,
        'organization_id': user['organization_id'],
        'owner_id': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.customers.insert_one(customer_doc)
    
    # Update client with customer_id
    await db.clients.update_one(
        {'id': client_id},
        {'$set': {'customer_id': customer_id}}
    )
    
    # Log activity
    await log_activity(
        org_id=user['organization_id'],
        user_id=user['id'],
        activity_type='conversion',
        subject="Lead converted to client",
        description=f"Lead \"{lead.get('name')}\" was converted to a client with {len(client_data.services)} service(s) worth ${total_value:,.2f}",
        lead_id=lead_id
    )
    
    # Notify the lead owner
    if lead.get('owner_id') and lead['owner_id'] != user['id']:
        await create_notification(
            org_id=user['organization_id'],
            user_id=lead['owner_id'],
            notification_type='lead_converted',
            title='Lead Converted to Client',
            message=f'Lead "{lead.get("name")}" has been converted to a client',
            link='/clients',
            entity_type='client',
            entity_id=client_id
        )
    
    client_doc.pop('_id', None)
    return {"message": "Lead converted to client successfully", "client": client_doc}

@api_router.post("/clients/{client_id}/services")
async def add_service_to_client(client_id: str, service: ServiceCreate, user: dict = Depends(get_current_user)):
    """Add a service to an existing client"""
    check_permission(user, Permission.MANAGE_CONTACTS)
    
    client = await db.clients.find_one(
        {'id': client_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    
    if not client:
        raise HTTPException(status_code=404, detail="Client not found")
    
    now = datetime.now(timezone.utc).isoformat()
    service_data = service.model_dump()
    service_data['id'] = str(uuid.uuid4())
    service_data['added_at'] = now
    
    # Add service and update total
    new_total = client.get('total_value', 0) + service.amount
    
    await db.clients.update_one(
        {'id': client_id},
        {
            '$push': {'services': service_data},
            '$set': {'total_value': new_total, 'updated_at': now}
        }
    )
    
    # Update lifecycle to repeat customer if they have multiple services
    services_count = len(client.get('services', [])) + 1
    if services_count > 1 and client.get('lead_id'):
        await db.leads.update_one(
            {'id': client['lead_id']},
            {'$set': {'lifecycle_stage': 'repeat_customer', 'updated_at': now}}
        )
    
    return {"message": "Service added successfully", "service": service_data}

@api_router.put("/leads/{lead_id}/pipeline-stage")
async def update_lead_pipeline_stage(lead_id: str, stage: str, user: dict = Depends(get_current_user)):
    """Update lead pipeline stage (after call completion)"""
    if Permission.MANAGE_ALL_LEADS.value not in user['permissions'] and Permission.MANAGE_OWN_LEADS.value not in user['permissions']:
        raise HTTPException(status_code=403, detail="Permission denied")
    
    valid_stages = [s.value for s in PipelineStage]
    if stage not in valid_stages:
        raise HTTPException(status_code=400, detail=f"Invalid stage. Valid stages: {valid_stages}")
    
    lead = await db.leads.find_one(
        {'id': lead_id, 'organization_id': user.get('organization_id')},
        {'_id': 0, 'pipeline_stage': 1, 'name': 1}
    )
    
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    old_stage = lead.get('pipeline_stage', 'new')
    now = datetime.now(timezone.utc).isoformat()
    
    # Update lifecycle based on pipeline stage
    lifecycle_mapping = {
        'new': 'lead',
        'contacted': 'ai_contacted',
        'no_answer': 'ai_contacted',
        'interested': 'interested',
        'follow_up': 'interested',
        'booked': 'opportunity',
        'won': 'customer',
        'lost': 'lead'
    }
    lifecycle_stage = lifecycle_mapping.get(stage, 'lead')
    
    await db.leads.update_one(
        {'id': lead_id},
        {'$set': {
            'pipeline_stage': stage,
            'lifecycle_stage': lifecycle_stage,
            'updated_at': now
        }}
    )
    
    # Log activity
    await log_activity(
        org_id=user['organization_id'],
        user_id=user['id'],
        activity_type='pipeline_change',
        subject=f"Pipeline stage changed: {old_stage} → {stage}",
        description=f"Lead moved from {old_stage} to {stage}",
        lead_id=lead_id
    )
    
    return {"message": f"Lead moved to {stage}", "pipeline_stage": stage, "lifecycle_stage": lifecycle_stage}

# ============= TASK ROUTES =============

@api_router.get("/tasks")
async def get_tasks(
    status: Optional[str] = None,
    payment_status: Optional[str] = None,
    assigned_to: Optional[str] = None,
    deal_id: Optional[str] = None,
    search: Optional[str] = None,
    page: int = 1,
    limit: int = 20,
    user: dict = Depends(get_current_user)
):
    """Get tasks for organization"""
    if not user.get('organization_id'):
        return {"items": [], "total": 0, "page": 1, "limit": limit, "total_pages": 0}
    
    page = max(1, page)
    limit = min(max(1, limit), 100)
    skip = (page - 1) * limit
    
    query = {'organization_id': user['organization_id']}
    
    if status:
        query['status'] = status
    if payment_status:
        query['payment_status'] = payment_status
    if assigned_to:
        query['assigned_to'] = assigned_to
    if deal_id:
        query['deal_id'] = deal_id
    if search:
        query['$or'] = [
            {'title': {'$regex': search, '$options': 'i'}},
            {'description': {'$regex': search, '$options': 'i'}}
        ]
    
    total = await db.tasks.count_documents(query)
    total_pages = (total + limit - 1) // limit
    
    tasks = await db.tasks.find(query, {'_id': 0}).sort('created_at', -1).skip(skip).limit(limit).to_list(limit)
    
    # Populate related data
    for task in tasks:
        task['reg_time'] = task.get('created_at')
        
        if task.get('assigned_to'):
            assignee = await db.users.find_one({'id': task['assigned_to']}, {'_id': 0, 'name': 1})
            task['assigned_to_name'] = assignee['name'] if assignee else None
        
        if task.get('lead_id'):
            lead = await db.leads.find_one({'id': task['lead_id']}, {'_id': 0, 'name': 1, 'company': 1, 'pic_name': 1})
            if lead:
                task['lead_name'] = lead.get('name')
                task['company_name'] = lead.get('company') or lead.get('name')
                task['pic_name'] = lead.get('pic_name')  # Don't fallback to name, keep it blank if not set
        
        if task.get('client_id'):
            client = await db.clients.find_one({'id': task['client_id']}, {'_id': 0, 'customer_name': 1, 'company': 1})
            if client:
                task['client_name'] = client.get('customer_name')
                if not task.get('company_name'):
                    task['company_name'] = client.get('company') or client.get('customer_name')
        
        if task.get('deal_id'):
            deal = await db.deals.find_one({'id': task['deal_id']}, {'_id': 0, 'title': 1, 'value': 1})
            if deal:
                task['deal_name'] = deal.get('title')
                task['deal_value'] = deal.get('value')
            
            # Fetch pipeline_status from lead-deal linkage if both lead_id and deal_id exist
            if task.get('lead_id'):
                linkage = await db.lead_deal_linkages.find_one({
                    'lead_id': task['lead_id'],
                    'deal_id': task['deal_id'],
                    'organization_id': user['organization_id']
                }, {'_id': 0, 'pipeline_status': 1})
                if linkage:
                    task['pipeline_status'] = linkage.get('pipeline_status')
    
    return {
        "items": tasks,
        "total": total,
        "page": page,
        "limit": limit,
        "total_pages": total_pages
    }

@api_router.post("/tasks")
async def create_task(task_data: TaskCreate, user: dict = Depends(get_current_user)):
    """Create a new task"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    now = datetime.now(timezone.utc).isoformat()
    task_id = str(uuid.uuid4())
    
    task_doc = {
        'id': task_id,
        'title': task_data.title,
        'description': task_data.description,
        'lead_id': task_data.lead_id,
        'client_id': task_data.client_id,
        'deal_id': task_data.deal_id,
        'assigned_to': task_data.assigned_to or user['id'],
        'due_date': task_data.due_date,
        'payment_status': task_data.payment_status,
        'payment_amount': task_data.payment_amount,
        'paid_amount': task_data.paid_amount or 0,
        'status': task_data.status or 'pending',
        'priority': task_data.priority,
        'calendar_event_id': None,
        'organization_id': user['organization_id'],
        'created_by': user['id'],
        'created_at': now,
        'updated_at': now
    }
    
    await db.tasks.insert_one(task_doc)
    
    # Notify assignee if different from creator
    if task_data.assigned_to and task_data.assigned_to != user['id']:
        await create_notification(
            org_id=user['organization_id'],
            user_id=task_data.assigned_to,
            notification_type='task_due',
            title='New Task Assigned',
            message=f'Task "{task_data.title}" has been assigned to you',
            link='/tasks',
            entity_type='task',
            entity_id=task_id
        )
    
    task_doc.pop('_id', None)
    return task_doc

@api_router.put("/tasks/{task_id}")
async def update_task(task_id: str, task_data: TaskUpdate, user: dict = Depends(get_current_user)):
    """Update a task"""
    task = await db.tasks.find_one(
        {'id': task_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    now = datetime.now(timezone.utc).isoformat()
    update_data = {k: v for k, v in task_data.model_dump().items() if v is not None}
    update_data['updated_at'] = now
    
    # Auto-update payment status based on paid_amount
    if 'paid_amount' in update_data and task.get('payment_amount'):
        paid = update_data['paid_amount']
        total = task['payment_amount']
        if paid >= total:
            update_data['payment_status'] = 'paid'
        elif paid > 0:
            update_data['payment_status'] = 'partially_paid'
        else:
            update_data['payment_status'] = 'unpaid'
    
    await db.tasks.update_one({'id': task_id}, {'$set': update_data})
    
    return {"message": "Task updated successfully"}

@api_router.delete("/tasks/{task_id}")
async def delete_task(task_id: str, user: dict = Depends(get_current_user)):
    """Delete a task"""
    result = await db.tasks.delete_one(
        {'id': task_id, 'organization_id': user.get('organization_id')}
    )
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Task not found")
    
    return {"message": "Task deleted successfully"}

# ============= LOOKUP DATA ROUTES =============

@api_router.get("/lookup/states")
async def get_states(user: dict = Depends(get_current_user)):
    """Get unique states from leads for filtering"""
    if not user.get('organization_id'):
        return {"states": []}
    
    pipeline = [
        {'$match': {'organization_id': user['organization_id'], 'state': {'$nin': [None, '']}}},
        {'$group': {'_id': '$state'}},
        {'$sort': {'_id': 1}}
    ]
    
    states = await db.leads.aggregate(pipeline).to_list(100)
    return {"states": [s['_id'] for s in states if s['_id']]}

@api_router.get("/lookup/sales-persons")
async def get_sales_persons(user: dict = Depends(get_current_user)):
    """Get users who can be assigned to tasks/leads"""
    if not user.get('organization_id'):
        return {"sales_persons": []}
    
    users = await db.users.find(
        {'organization_id': user['organization_id'], 'is_active': True},
        {'_id': 0, 'id': 1, 'name': 1, 'role': 1}
    ).to_list(100)
    
    return {"sales_persons": users}

@api_router.get("/lookup/companies")
async def get_companies_for_linking(
    search: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get leads and customers for linking to deals"""
    if not user.get('organization_id'):
        return {"companies": []}
    
    query = {'organization_id': user['organization_id']}
    if search:
        query['$or'] = [
            {'name': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}},
            {'pic_name': {'$regex': search, '$options': 'i'}}
        ]
    
    # Get leads
    leads = await db.leads.find(
        query,
        {'_id': 0, 'id': 1, 'name': 1, 'company': 1, 'pic_name': 1, 'phone': 1, 'city': 1, 'state': 1}
    ).limit(20).to_list(20)
    
    companies = []
    for lead in leads:
        companies.append({
            'id': lead['id'],
            'name': lead.get('company') or lead.get('name'),
            'pic_name': lead.get('pic_name') or lead.get('name'),
            'phone': lead.get('phone'),
            'location': f"{lead.get('city', '')}, {lead.get('state', '')}".strip(', '),
            'entity_type': 'lead'
        })
    
    # Get customers
    customer_query = {'organization_id': user['organization_id']}
    if search:
        customer_query['$or'] = [
            {'first_name': {'$regex': search, '$options': 'i'}},
            {'last_name': {'$regex': search, '$options': 'i'}},
            {'company': {'$regex': search, '$options': 'i'}}
        ]
    
    customers = await db.customers.find(
        customer_query,
        {'_id': 0, 'id': 1, 'first_name': 1, 'last_name': 1, 'company': 1, 'phone': 1, 'city': 1, 'state': 1}
    ).limit(20).to_list(20)
    
    for customer in customers:
        companies.append({
            'id': customer['id'],
            'name': customer.get('company') or f"{customer.get('first_name', '')} {customer.get('last_name', '')}".strip(),
            'pic_name': f"{customer.get('first_name', '')} {customer.get('last_name', '')}".strip(),
            'phone': customer.get('phone'),
            'location': f"{customer.get('city', '')}, {customer.get('state', '')}".strip(', '),
            'entity_type': 'customer'
        })
    
    return {"companies": companies}

# ============= ORGANIZATION SETTINGS ROUTES =============

@api_router.get("/organization-settings")
async def get_organization_settings(user: dict = Depends(get_current_user)):
    """Get organization settings including currency and integrations"""
    check_permission(user, Permission.VIEW_ORGANIZATION)
    
    if not user.get('organization_id'):
        return {
            "organization_id": None,
            "currency": "USD",
            "currency_symbol": "$",
            "google_calendar_enabled": False,
            "google_calendar_connected": False,
            "twilio_enabled": False,
            "twilio_connected": False,
            "twilio_whatsapp_number": None
        }
    
    settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0, 'google_calendar_client_secret': 0, 'twilio_auth_token': 0}  # Don't return secrets
    )
    
    if not settings:
        return {
            "organization_id": user['organization_id'],
            "currency": "USD",
            "currency_symbol": "$",
            "google_calendar_enabled": False,
            "google_calendar_connected": False,
            "twilio_enabled": False,
            "twilio_connected": False,
            "twilio_whatsapp_number": None
        }
    
    # Check if Google Calendar credentials are configured
    full_settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    settings['google_calendar_connected'] = bool(
        full_settings.get('google_calendar_client_id') and 
        full_settings.get('google_calendar_client_secret') and
        full_settings.get('google_calendar_access_token')
    )
    
    # Check if Twilio credentials are configured
    settings['twilio_connected'] = bool(
        full_settings.get('twilio_account_sid') and 
        full_settings.get('twilio_auth_token') and
        full_settings.get('twilio_whatsapp_number')
    )
    settings['twilio_enabled'] = full_settings.get('twilio_enabled', False)
    settings['twilio_whatsapp_number'] = full_settings.get('twilio_whatsapp_number')
    
    return settings

@api_router.put("/organization-settings")
async def update_organization_settings(settings: OrganizationSettingsUpdate, user: dict = Depends(get_current_user)):
    """Update organization settings"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    now = datetime.now(timezone.utc).isoformat()
    
    update_data = {k: v for k, v in settings.model_dump().items() if v is not None}
    update_data['organization_id'] = user['organization_id']
    update_data['updated_at'] = now
    
    await db.organization_settings.update_one(
        {'organization_id': user['organization_id']},
        {'$set': update_data},
        upsert=True
    )
    
    return {"message": "Organization settings updated successfully"}

@api_router.post("/twilio/test-connection")
async def test_twilio_connection(user: dict = Depends(get_current_user)):
    """Test Twilio WhatsApp connection with saved credentials"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not settings or not settings.get('twilio_account_sid') or not settings.get('twilio_auth_token'):
        raise HTTPException(status_code=400, detail="Twilio credentials not configured. Please save your Account SID and Auth Token first.")
    
    try:
        twilio_client = TwilioClient(settings['twilio_account_sid'], settings['twilio_auth_token'])
        # Verify credentials by fetching account info
        account = twilio_client.api.accounts(settings['twilio_account_sid']).fetch()
        
        return {
            "success": True,
            "message": "Twilio connection successful!",
            "account_name": account.friendly_name,
            "account_status": account.status
        }
    except TwilioRestException as e:
        logger.error(f"Twilio connection test failed: {e.msg}")
        raise HTTPException(status_code=400, detail=f"Twilio connection failed: {e.msg}")
    except Exception as e:
        logger.error(f"Twilio connection test error: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to connect to Twilio")

@api_router.get("/currencies")
async def get_available_currencies():
    """Get list of available currencies"""
    return {
        "currencies": [
            {"code": "USD", "symbol": "$", "name": "US Dollar"},
            {"code": "MYR", "symbol": "RM", "name": "Malaysian Ringgit"},
            {"code": "EUR", "symbol": "€", "name": "Euro"},
            {"code": "GBP", "symbol": "£", "name": "British Pound"},
            {"code": "SGD", "symbol": "S$", "name": "Singapore Dollar"},
            {"code": "AUD", "symbol": "A$", "name": "Australian Dollar"},
            {"code": "JPY", "symbol": "¥", "name": "Japanese Yen"},
            {"code": "CNY", "symbol": "¥", "name": "Chinese Yuan"},
            {"code": "INR", "symbol": "₹", "name": "Indian Rupee"},
            {"code": "THB", "symbol": "฿", "name": "Thai Baht"},
            {"code": "IDR", "symbol": "Rp", "name": "Indonesian Rupiah"},
            {"code": "PHP", "symbol": "₱", "name": "Philippine Peso"},
            {"code": "VND", "symbol": "₫", "name": "Vietnamese Dong"},
            {"code": "KRW", "symbol": "₩", "name": "South Korean Won"},
            {"code": "AED", "symbol": "د.إ", "name": "UAE Dirham"},
            {"code": "SAR", "symbol": "﷼", "name": "Saudi Riyal"}
        ]
    }

# ============= GOOGLE CALENDAR ROUTES =============

@api_router.get("/google-calendar/auth-url")
async def get_google_calendar_auth_url(user: dict = Depends(get_current_user)):
    """Get Google OAuth authorization URL"""
    check_permission(user, Permission.MANAGE_ORGANIZATION)
    
    settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not settings or not settings.get('google_calendar_client_id'):
        raise HTTPException(status_code=400, detail="Google Calendar credentials not configured")
    
    client_id = settings['google_calendar_client_id']
    redirect_uri = os.environ.get('REACT_APP_BACKEND_URL', '') + '/api/google-calendar/callback'
    
    scope = 'https://www.googleapis.com/auth/calendar https://www.googleapis.com/auth/calendar.events'
    
    auth_url = (
        f"https://accounts.google.com/o/oauth2/v2/auth?"
        f"client_id={client_id}&"
        f"redirect_uri={redirect_uri}&"
        f"response_type=code&"
        f"scope={scope}&"
        f"access_type=offline&"
        f"prompt=consent&"
        f"state={user['organization_id']}"
    )
    
    return {"auth_url": auth_url}

@api_router.get("/google-calendar/callback")
async def google_calendar_callback(code: str, state: str):
    """Handle Google OAuth callback"""
    org_id = state
    
    settings = await db.organization_settings.find_one(
        {'organization_id': org_id},
        {'_id': 0}
    )
    
    if not settings:
        return {"error": "Organization settings not found"}
    
    client_id = settings.get('google_calendar_client_id')
    client_secret = settings.get('google_calendar_client_secret')
    redirect_uri = os.environ.get('REACT_APP_BACKEND_URL', '') + '/api/google-calendar/callback'
    
    # Exchange code for tokens
    import httpx
    async with httpx.AsyncClient() as client:
        response = await client.post(
            'https://oauth2.googleapis.com/token',
            data={
                'client_id': client_id,
                'client_secret': client_secret,
                'code': code,
                'grant_type': 'authorization_code',
                'redirect_uri': redirect_uri
            }
        )
        
        if response.status_code != 200:
            return {"error": "Failed to exchange code for tokens"}
        
        tokens = response.json()
    
    # Store tokens
    now = datetime.now(timezone.utc).isoformat()
    await db.organization_settings.update_one(
        {'organization_id': org_id},
        {'$set': {
            'google_calendar_access_token': tokens.get('access_token'),
            'google_calendar_refresh_token': tokens.get('refresh_token'),
            'google_calendar_token_expiry': tokens.get('expires_in'),
            'google_calendar_connected_at': now,
            'updated_at': now
        }}
    )
    
    # Redirect to settings page
    frontend_url = os.environ.get('REACT_APP_BACKEND_URL', '').replace('/api', '')
    from fastapi.responses import RedirectResponse
    return RedirectResponse(url=f"{frontend_url}/settings?google_calendar=connected")

@api_router.post("/google-calendar/sync-task/{task_id}")
async def sync_task_to_calendar(task_id: str, user: dict = Depends(get_current_user)):
    """Sync a task to Google Calendar"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    # Get task
    task = await db.tasks.find_one(
        {'id': task_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    # Get Google Calendar settings
    settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not settings or not settings.get('google_calendar_access_token'):
        raise HTTPException(status_code=400, detail="Google Calendar not connected")
    
    access_token = settings['google_calendar_access_token']
    
    # Create calendar event
    import httpx
    event = {
        'summary': task['title'],
        'description': task.get('description', ''),
        'start': {
            'dateTime': task.get('due_date') or datetime.now(timezone.utc).isoformat(),
            'timeZone': 'UTC'
        },
        'end': {
            'dateTime': task.get('due_date') or datetime.now(timezone.utc).isoformat(),
            'timeZone': 'UTC'
        }
    }
    
    async with httpx.AsyncClient() as client:
        response = await client.post(
            'https://www.googleapis.com/calendar/v3/calendars/primary/events',
            headers={'Authorization': f'Bearer {access_token}'},
            json=event
        )
        
        if response.status_code not in [200, 201]:
            raise HTTPException(status_code=400, detail="Failed to create calendar event")
        
        calendar_event = response.json()
    
    # Update task with calendar event ID
    await db.tasks.update_one(
        {'id': task_id},
        {'$set': {'calendar_event_id': calendar_event.get('id')}}
    )
    
    return {"message": "Task synced to Google Calendar", "event_id": calendar_event.get('id')}

# ============= CALENDAR EVENTS ROUTES =============

class CalendarEventCreate(BaseModel):
    title: str
    description: Optional[str] = None
    date: str  # YYYY-MM-DD
    start_time: Optional[str] = "09:00"
    end_time: Optional[str] = "10:00"
    location: Optional[str] = None
    color: Optional[str] = "#A0C4FF"
    all_day: Optional[bool] = False
    attendees: List[str] = []

class CalendarEventUpdate(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    date: Optional[str] = None
    start_time: Optional[str] = None
    end_time: Optional[str] = None
    location: Optional[str] = None
    color: Optional[str] = None
    all_day: Optional[bool] = None

@api_router.get("/calendar/events")
async def get_calendar_events(
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get calendar events for the organization"""
    if not user.get('organization_id'):
        return {"events": []}
    
    query = {'organization_id': user['organization_id']}
    
    if start_date and end_date:
        query['date'] = {'$gte': start_date, '$lte': end_date}
    
    events = await db.calendar_events.find(query, {'_id': 0}).sort('date', 1).to_list(500)
    
    return {"events": events}

@api_router.post("/calendar/events")
async def create_calendar_event(event_data: CalendarEventCreate, user: dict = Depends(get_current_user)):
    """Create a new calendar event"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    event_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    event_doc = {
        'id': event_id,
        **event_data.model_dump(),
        'organization_id': user['organization_id'],
        'created_by': user['id'],
        'created_at': now,
        'updated_at': now,
        'google_event_id': None
    }
    
    await db.calendar_events.insert_one(event_doc)
    event_doc.pop('_id', None)
    
    return event_doc

@api_router.get("/calendar/events/{event_id}")
async def get_calendar_event(event_id: str, user: dict = Depends(get_current_user)):
    """Get a specific calendar event"""
    event = await db.calendar_events.find_one(
        {'id': event_id, 'organization_id': user.get('organization_id')},
        {'_id': 0}
    )
    
    if not event:
        raise HTTPException(status_code=404, detail="Event not found")
    
    return event

@api_router.put("/calendar/events/{event_id}")
async def update_calendar_event(
    event_id: str, 
    event_data: CalendarEventUpdate, 
    user: dict = Depends(get_current_user)
):
    """Update a calendar event"""
    update_data = {k: v for k, v in event_data.model_dump().items() if v is not None}
    update_data['updated_at'] = datetime.now(timezone.utc).isoformat()
    
    result = await db.calendar_events.find_one_and_update(
        {'id': event_id, 'organization_id': user.get('organization_id')},
        {'$set': update_data},
        return_document=True,
        projection={'_id': 0}
    )
    
    if not result:
        raise HTTPException(status_code=404, detail="Event not found")
    
    return result

@api_router.delete("/calendar/events/{event_id}")
async def delete_calendar_event(event_id: str, user: dict = Depends(get_current_user)):
    """Delete a calendar event"""
    result = await db.calendar_events.delete_one(
        {'id': event_id, 'organization_id': user.get('organization_id')}
    )
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Event not found")
    
    return {"message": "Event deleted successfully"}

# ============= WHATSAPP MESSAGE ROUTES =============

class WhatsAppMessageSend(BaseModel):
    contact_id: str
    message: str
    phone: Optional[str] = None

@api_router.get("/whatsapp/messages/{contact_id}")
async def get_whatsapp_messages(contact_id: str, user: dict = Depends(get_current_user)):
    """Get WhatsApp messages for a contact"""
    if not user.get('organization_id'):
        return {"messages": []}
    
    messages = await db.whatsapp_messages.find(
        {'organization_id': user['organization_id'], 'contact_id': contact_id},
        {'_id': 0}
    ).sort('timestamp', 1).to_list(100)
    
    return {"messages": messages}

@api_router.post("/whatsapp/send")
async def send_whatsapp_message(message_data: WhatsAppMessageSend, user: dict = Depends(get_current_user)):
    """Send a WhatsApp message via Twilio (if configured) or store locally"""
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first. Go to Organization settings to get started.")
    
    message_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    # Check if Twilio is configured for this organization
    settings = await db.organization_settings.find_one(
        {'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    twilio_message_sid = None
    twilio_status = None
    status = 'sent'  # Default for local storage
    
    if settings and settings.get('twilio_enabled') and settings.get('twilio_account_sid') and settings.get('twilio_auth_token') and settings.get('twilio_whatsapp_number'):
        try:
            # Send via Twilio
            twilio_client = TwilioClient(settings['twilio_account_sid'], settings['twilio_auth_token'])
            
            # Format phone number for WhatsApp
            to_number = message_data.phone
            if not to_number.startswith('whatsapp:'):
                to_number = f"whatsapp:{to_number}" if to_number.startswith('+') else f"whatsapp:+{to_number}"
            
            from_number = settings['twilio_whatsapp_number']
            if not from_number.startswith('whatsapp:'):
                from_number = f"whatsapp:{from_number}"
            
            message = twilio_client.messages.create(
                body=message_data.message,
                from_=from_number,
                to=to_number
            )
            
            twilio_message_sid = message.sid
            twilio_status = message.status
            status = 'delivered' if message.status in ['delivered', 'sent', 'queued'] else message.status
            logger.info(f"WhatsApp message sent via Twilio: SID={message.sid}, Status={message.status}")
            
        except TwilioRestException as e:
            logger.error(f"Twilio API error: {e.msg}")
            raise HTTPException(status_code=400, detail=f"Failed to send WhatsApp message: {e.msg}")
        except Exception as e:
            logger.error(f"Error sending WhatsApp via Twilio: {str(e)}")
            raise HTTPException(status_code=500, detail="Failed to send WhatsApp message")
    
    message_doc = {
        'id': message_id,
        'contact_id': message_data.contact_id,
        'phone': message_data.phone,
        'text': message_data.message,
        'sender': 'me',
        'timestamp': now,
        'status': status,
        'twilio_message_sid': twilio_message_sid,
        'twilio_status': twilio_status,
        'organization_id': user['organization_id'],
        'sent_by': user['id']
    }
    
    await db.whatsapp_messages.insert_one(message_doc)
    message_doc.pop('_id', None)
    
    # Log activity
    await db.activities.insert_one({
        'id': str(uuid.uuid4()),
        'organization_id': user['organization_id'],
        'entity_type': 'lead',
        'entity_id': message_data.contact_id,
        'action': 'whatsapp_sent',
        'title': 'WhatsApp Message Sent',
        'description': f'Sent: {message_data.message[:100]}...' if len(message_data.message) > 100 else f'Sent: {message_data.message}',
        'user_id': user['id'],
        'user_name': user['name'],
        'created_at': now
    })
    
    return {"success": True, "message": message_doc, "twilio_sid": twilio_message_sid}

# ============= BATCH AI CALLING & MESSAGING =============

class BatchAiCallRequest(BaseModel):
    lead_ids: List[str]
    agent_name: str
    call_purpose: Optional[str] = "follow_up"  # follow_up, appointment, qualification, custom

class BatchWhatsAppRequest(BaseModel):
    lead_ids: List[str]
    message_type: Optional[str] = "follow_up"  # follow_up, appointment, introduction, thank_you
    custom_message: Optional[str] = None

class AiWhatsAppMessageRequest(BaseModel):
    lead_id: str
    message_type: Optional[str] = "follow_up"
    custom_message: Optional[str] = None

@api_router.post("/ai-calls/batch")
async def batch_ai_call_endpoint(request: BatchAiCallRequest, user: dict = Depends(get_current_user)):
    """Initiate batch AI calls to multiple leads"""
    from ai_services import batch_ai_call
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    # Get leads
    leads = await db.leads.find(
        {'id': {'$in': request.lead_ids}, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(100)
    
    if not leads:
        raise HTTPException(status_code=404, detail="No leads found")
    
    # Get agent configuration
    agent = await db.ai_agents.find_one(
        {'organization_id': user['organization_id'], 'name': request.agent_name},
        {'_id': 0}
    )
    
    # Execute batch calls
    result = await batch_ai_call(
        leads=leads,
        agent_config=agent or {'name': request.agent_name},
        call_purpose=request.call_purpose
    )
    
    # Log batch activity
    now = datetime.now(timezone.utc).isoformat()
    await db.activities.insert_one({
        'id': str(uuid.uuid4()),
        'organization_id': user['organization_id'],
        'entity_type': 'batch',
        'entity_id': 'batch_call',
        'action': 'batch_ai_call',
        'title': f'Batch AI Call - {len(leads)} leads',
        'description': f'Initiated batch AI calls with {request.agent_name}. Success: {result["successful"]}, Failed: {result["failed"]}',
        'user_id': user['id'],
        'user_name': user['name'],
        'metadata': {
            'total': result['total'],
            'successful': result['successful'],
            'failed': result['failed'],
            'agent_name': request.agent_name,
            'call_purpose': request.call_purpose
        },
        'created_at': now
    })
    
    return result

@api_router.post("/whatsapp/batch")
async def batch_whatsapp_endpoint(request: BatchWhatsAppRequest, user: dict = Depends(get_current_user)):
    """Send batch WhatsApp messages to multiple leads"""
    from ai_services import batch_whatsapp_message
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    # Get leads
    leads = await db.leads.find(
        {'id': {'$in': request.lead_ids}, 'organization_id': user['organization_id']},
        {'_id': 0}
    ).to_list(100)
    
    if not leads:
        raise HTTPException(status_code=404, detail="No leads found")
    
    # Execute batch messaging
    result = await batch_whatsapp_message(
        leads=leads,
        message_type=request.message_type,
        custom_message=request.custom_message
    )
    
    # Log batch activity
    now = datetime.now(timezone.utc).isoformat()
    await db.activities.insert_one({
        'id': str(uuid.uuid4()),
        'organization_id': user['organization_id'],
        'entity_type': 'batch',
        'entity_id': 'batch_whatsapp',
        'action': 'batch_whatsapp',
        'title': f'Batch WhatsApp - {len(leads)} leads',
        'description': f'Sent batch WhatsApp messages. Success: {result["successful"]}, Failed: {result["failed"]}',
        'user_id': user['id'],
        'user_name': user['name'],
        'metadata': {
            'total': result['total'],
            'successful': result['successful'],
            'failed': result['failed'],
            'message_type': request.message_type
        },
        'created_at': now
    })
    
    return result

@api_router.post("/whatsapp/ai-send")
async def send_ai_whatsapp_message_endpoint(request: AiWhatsAppMessageRequest, user: dict = Depends(get_current_user)):
    """Send an AI-powered personalized WhatsApp message to a lead"""
    from ai_services import send_ai_whatsapp_message
    
    if not user.get('organization_id'):
        raise HTTPException(status_code=400, detail="Please create or join an organization first")
    
    # Get lead
    lead = await db.leads.find_one(
        {'id': request.lead_id, 'organization_id': user['organization_id']},
        {'_id': 0}
    )
    
    if not lead:
        raise HTTPException(status_code=404, detail="Lead not found")
    
    phone = lead.get('phone') or lead.get('office_number')
    if not phone:
        raise HTTPException(status_code=400, detail="No phone number available for this lead")
    
    # Get deal info if available
    deal = None
    linkage = await db.lead_deal_linkages.find_one({'lead_id': request.lead_id}, {'_id': 0})
    if linkage:
        deal = await db.deals.find_one({'id': linkage['deal_id']}, {'_id': 0})
    
    # Send AI WhatsApp message
    result = await send_ai_whatsapp_message(
        recipient_phone=phone,
        lead_info=lead,
        deal_info=deal,
        message_type=request.message_type,
        custom_message=request.custom_message
    )
    
    # Store message in database
    if result.get('success'):
        now = datetime.now(timezone.utc).isoformat()
        message_doc = {
            'id': str(uuid.uuid4()),
            'organization_id': user['organization_id'],
            'contact_id': request.lead_id,
            'direction': 'outgoing',
            'sender': 'AI',
            'content': result.get('body', ''),
            'message_type': request.message_type,
            'phone': phone,
            'twilio_sid': result.get('message_sid'),
            'status': result.get('status', 'sent'),
            'timestamp': now
        }
        await db.whatsapp_messages.insert_one(message_doc)
        
        # Log activity
        await db.activities.insert_one({
            'id': str(uuid.uuid4()),
            'organization_id': user['organization_id'],
            'entity_type': 'lead',
            'entity_id': request.lead_id,
            'action': 'ai_whatsapp_sent',
            'title': 'AI WhatsApp Message Sent',
            'description': f'AI {request.message_type} message sent to {lead.get("name") or lead.get("company")}',
            'user_id': user['id'],
            'user_name': user['name'],
            'created_at': now
        })
    
    return result

@api_router.get("/elevenlabs/voices")
async def get_elevenlabs_voices(user: dict = Depends(get_current_user)):
    """Get available ElevenLabs voices for AI calling"""
    from ai_services import get_available_voices
    
    voices = await get_available_voices()
    return {"voices": voices}

# ============= ROOT ROUTE =============

@api_router.get("/")
async def root():
    return {"message": "SalesHub Platform API", "version": "2.0.0", "features": ["Multi-tenant", "RBAC", "CRM", "Sales"]}

# Include router and middleware
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
